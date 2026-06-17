"""JavaBrain PPT 工具库 —— 跨版本(v1/v2/v3)可复用的工具函数。

分组:
A. 文本/形状:   add_text, set_text, add_rect, add_picture_with_alpha
B. 布局/装饰:   set_solid_bg, grid_bg, hud_corner, section_label,
                corner_badge, add_page_number, apply_chrome
C. 复合元素:    kill_box, node_block, card, terminal_box, big_num,
                iso_lego, arrow_line
D. 文件 I/O:    replace_pptx_media, find_dominant_bg, cover_region,
                qa_image_check, new_presentation
E. 万相 API:    wanx_generate, wanx_upload, wanx_edit, wanx_download

排版常量(顶部集中 magic numbers):TEXT_MARGIN_*, STUD_*, QA_*

注:动画逻辑见 scripts/pptx_anim.py(独立模块,可选)

依赖:python-pptx, lxml, pillow。
"""
from __future__ import annotations

import json
import os
import sys
import urllib.error
import urllib.request
import zipfile
from collections import Counter
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ====================================================================
# 设计 Token(暗色 HUD 主题 —— 跨版本统一)
# ====================================================================

# 字体
FONT_CN = "Microsoft YaHei"
FONT_EN = "Inter"
FONT_MONO = "JetBrains Mono"

# 暗色背景
BG_DEEP = RGBColor(0x0A, 0x0E, 0x1A)       # 主背景(深色 HUD)
BG_PANEL = RGBColor(0x10, 0x17, 0x2A)      # 次级面板
GRID_LINE = RGBColor(0x1A, 0x22, 0x38)     # 网格线 / 卡片描边
DIVIDER = RGBColor(0x1A, 0x22, 0x38)

# 文字
TEXT_PRIMARY = RGBColor(0xE4, 0xE7, 0xF1)   # 主文字
TEXT_SECONDARY = RGBColor(0x7C, 0x8A, 0xA8) # 次级文字
TEXT_DIM = RGBColor(0x4A, 0x55, 0x6E)       # 注脚

# 语义主色(中性命名,不绑定项目)
PRIMARY = RGBColor(0x00, 0xD9, 0xFF)       # 主色
ACCENT = RGBColor(0xA7, 0x8B, 0xFA)        # 强调色
SUCCESS = RGBColor(0x00, 0xFF, 0x9C)       # 成功 / 反衬
DANGER = RGBColor(0xFF, 0x4D, 0x6D)        # 警示 / 痛点
WARN = RGBColor(0xFA, 0xCC, 0x15)          # 警示金 / 重点

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# ====================================================================
# 排版常量(集中 magic numbers)
# ====================================================================

# textbox 边距
TEXT_MARGIN_LR_IN = 0.05
TEXT_MARGIN_TB_IN = 0.02

# 2.5D 乐高积木参数(iso_lego)
STUD_RATIO = 0.16            # 圆钉半径 / 边长
STUD_Y_OFFSET_RATIO = 0.5    # 圆钉 Y 偏移(向上突出积木顶部的比例)
STUD_X_POSITIONS = (0.18, 0.40, 0.60, 0.82)  # 4 个圆钉 X 位置比例

# QA 阈值(qa_image_check)
QA_MIN_SHORT_SIDE = 720
QA_MIN_FILE_BYTES = 200_000
QA_MAX_FILE_BYTES = 2_000_000
QA_MIN_STD = 30
QA_MIN_RATIO = 0.5
QA_MAX_RATIO = 3.0


# ====================================================================
# A. 文本 / 形状
# ====================================================================

def add_text(slide, x_in, y_in, w_in, h_in, text, *,
             font=FONT_MONO, size=18, color=TEXT_PRIMARY,
             bold=False, italic=False, align=PP_ALIGN.CENTER):
    """添加一个 textbox 并写文字,返回 textbox(用于 add_anim)。"""
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(h_in))
    set_text(tb, text, font=font, size=size, color=color,
             bold=bold, italic=italic, align=align)
    return tb


def set_text(shape, text, *, font=FONT_MONO, size=18, color=TEXT_PRIMARY,
             bold=False, italic=False, align=PP_ALIGN.CENTER):
    """在已存在的 shape/textbox 上写文字 + 字体样式。

    字体约定:英文/数字走 ``font``,中文通过 a:ea 走 FONT_CN。
    """
    tf = shape.text_frame
    tf.text = text
    tf.margin_left = Inches(TEXT_MARGIN_LR_IN)
    tf.margin_right = Inches(TEXT_MARGIN_LR_IN)
    tf.margin_top = Inches(TEXT_MARGIN_TB_IN)
    tf.margin_bottom = Inches(TEXT_MARGIN_TB_IN)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    if not p.runs:
        return
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # 让 eastAsia 也走中文字体
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT_CN)


def add_rect(slide, x_in, y_in, w_in, h_in, *, fill=None, line_color=None,
             line_width_pt=0, shape=MSO_SHAPE.RECTANGLE):
    """添加一个矩形形状,可选填充和描边。返回 shape。"""
    s = slide.shapes.add_shape(shape, Inches(x_in), Inches(y_in),
                                Inches(w_in), Inches(h_in))
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width_pt)
    return s


def add_picture_with_alpha(slide, image_path, x, y, w, h, alpha_pct):
    """Add picture with overall alpha (0-100) via blip alphaModFix."""
    from PIL import Image  # 仅在用到时导入
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y),
                                    Inches(w), Inches(h))
    blip = pic._element.blipFill.blip
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    alpha = blip.find("a:alphaModFix", nsmap)
    if alpha is None:
        alpha = etree.SubElement(blip, qn("a:alphaModFix"))
    alpha.set("amt", str(int(alpha_pct * 1000)))
    return pic


# ====================================================================
# C. 布局 / 装饰
# ====================================================================

def set_solid_bg(slide, color) -> None:
    """设置 slide 背景为纯色。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def grid_bg(slide, cols=12, rows=8, color=GRID_LINE, weight_pt=0.5,
            slide_w_in=13.333, slide_h_in=7.5) -> None:
    """画 cols × rows 细网格背景(覆盖整张 slide)。"""
    for c in range(cols + 1):
        x = c * slide_w_in / cols
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x), Inches(0), Inches(x), Inches(slide_h_in))
        ln.line.color.rgb = color
        ln.line.width = Pt(weight_pt)
    for r in range(rows + 1):
        y = r * slide_h_in / rows
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0), Inches(y), Inches(slide_w_in), Inches(y))
        ln.line.color.rgb = color
        ln.line.width = Pt(weight_pt)


def hud_corner(slide, page_num, total=12, *, header='', pager='[ {n:02d} / {t} ]'):
    """HUD 角标(左上 + 右下页码)。header 为空时不显示左上文字。

    Args:
        header: 左上角文字(项目可填入自己的版本号 / 状态条 / 标语)
        pager: 右下角页码模板,用 {n} {t} 占位
    """
    if header:
        add_text(slide, 0.3, 0.18, 9.0, 0.3, header,
                 font=FONT_MONO, size=10, color=TEXT_SECONDARY,
                 align=PP_ALIGN.LEFT)
    add_text(slide, 10.5, 7.15, 2.5, 0.3,
             pager.format(n=page_num, t=total),
             font=FONT_MONO, size=10, color=TEXT_SECONDARY,
             align=PP_ALIGN.RIGHT)


def section_label(slide, text, *,
                  size=11, color=TEXT_SECONDARY,
                  prefix='', font=FONT_MONO, bold=False):
    """顶部章节标。prefix 用于装饰前缀(如 '── ')。

    Args:
        prefix: 文本前拼接的装饰字符串(如 '── ')
        size/color/font/bold: 透传给内部 add_text
    """
    add_text(slide, 0.5, 0.55, 12.333, 0.3,
             f"{prefix}{text}", font=font, size=size,
             color=color, bold=bold, align=PP_ALIGN.LEFT)




def corner_badge(slide, x, y, w, h, text, *,
                 fill=BG_PANEL, color=SUCCESS):
    """角落徽章(SUCCESS 边框 + 深底,标注场景/补充信息)。"""
    rect = add_rect(slide, x, y, w, h,
                     fill=fill, line_color=color, line_width_pt=1.5,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x, y, w, h, text,
             font=FONT_CN, size=11, color=color, bold=True)
    return rect


def add_page_number(slide, current, total=10, color=TEXT_SECONDARY, size=10):
    """右下角 'XX/total' 页码。"""
    add_text(slide, 11.6, 7.05, 1.4, 0.3, f"{current:02d} / {total}",
             font=FONT_MONO, size=size, color=color, align=PP_ALIGN.RIGHT)


def apply_hud_chrome(slide, page_num, total, section_text) -> None:
    """v3 标准化 slide 装饰四件套:背景 + 网格 + HUD 角标 + 章节标。

    等价于连续调用:
        set_solid_bg(slide, BG_DEEP)
        grid_bg(slide)
        hud_corner(slide, page_num, total=total)
        section_label(slide, section_text)

    适合 v3 / 后续版本每页都套这层 chrome。
    """
    set_solid_bg(slide, BG_DEEP)
    grid_bg(slide)
    hud_corner(slide, page_num, total=total)
    section_label(slide, section_text)


# ====================================================================
# D. 复合元素
# ====================================================================

def kill_box(slide, x_in, y_in, w_in, h_in, text, *,
             size=20, color=SUCCESS, fill=BG_PANEL):
    """金句框(SUCCESS 边框 2.5pt,深底)。"""
    box = add_rect(slide, x_in, y_in, w_in, h_in,
                    fill=fill, line_color=color, line_width_pt=2.5)
    add_text(slide, x_in, y_in, w_in, h_in, text,
             font=FONT_MONO, size=size, color=color, bold=True)
    return box


def node_block(slide, x_in, y_in, w_in, h_in, label, sublabel=None, *,
               color=PRIMARY, star=False, size_label=18, size_sub=11):
    """节点方块(架构图用)。返回主形状。"""
    rect = add_rect(slide, x_in, y_in, w_in, h_in,
                     fill=BG_PANEL, line_color=color, line_width_pt=2,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if star:
        add_text(slide, x_in + w_in - 0.4, y_in + 0.05, 0.35, 0.3, "★",
                 font=FONT_MONO, size=12, color=SUCCESS,
                 align=PP_ALIGN.RIGHT)
    add_text(slide, x_in, y_in + 0.15, w_in, 0.4, label,
             font=FONT_MONO, size=size_label, color=color, bold=True)
    if sublabel:
        add_text(slide, x_in, y_in + 0.55, w_in, 0.3, sublabel,
                 font=FONT_MONO, size=size_sub, color=TEXT_SECONDARY)
    return rect


def card(slide, x_in, y_in, w_in, h_in, *, fill=BG_PANEL,
         border=DIVIDER, border_pt=1, star=False):
    """通用卡片(功能模块用)。返回主形状。"""
    rect = add_rect(slide, x_in, y_in, w_in, h_in,
                     fill=fill, line_color=border, line_width_pt=border_pt,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if star:
        add_text(slide, x_in + w_in - 0.4, y_in + 0.05, 0.35, 0.25, "★",
                 font=FONT_MONO, size=12, color=SUCCESS,
                 align=PP_ALIGN.RIGHT)
    return rect


def terminal_box(slide, x_in, y_in, w_in, h_in, lines, *,
                 title=None, size=12):
    """终端框(深底 + 等宽字体,多行)。lines 是字符串列表。

    行颜色启发式:
      $   → SUCCESS(命令)
      >   → PRIMARY(输出)
      ✓/READY → SUCCESS 加粗
      ★  → SUCCESS 加粗
      #   → TEXT_SECONDARY(注释)
    """
    rect = add_rect(slide, x_in, y_in, w_in, h_in,
                     fill=BG_PANEL, line_color=DIVIDER, line_width_pt=1,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    line_h = 0.32
    start_y = y_in + (0.4 if title else 0.15)
    if title:
        add_text(slide, x_in + 0.15, y_in + 0.05, w_in - 0.3, 0.3, title,
                 font=FONT_MONO, size=11, color=TEXT_SECONDARY,
                 align=PP_ALIGN.LEFT)
    for i, line in enumerate(lines):
        color, bold = _terminal_line_style(line)
        add_text(slide, x_in + 0.15, start_y + i * line_h,
                 w_in - 0.3, line_h, line,
                 font=FONT_MONO, size=size, color=color, bold=bold,
                 align=PP_ALIGN.LEFT)
    return rect


def _terminal_line_style(line: str) -> tuple[RGBColor, bool]:
    """终端行颜色启发式(给 terminal_box 用)。"""
    if line.startswith("$"):
        return SUCCESS, False
    if line.startswith(">"):
        return PRIMARY, False
    if line.startswith("✓") or line.startswith("READY"):
        return SUCCESS, True
    if "★" in line:
        return SUCCESS, True
    if line.startswith("#"):
        return TEXT_SECONDARY, False
    return TEXT_PRIMARY, False


def big_num(slide, x_in, y_in, w_in, h_in, num, unit, label, *,
            color=PRIMARY, num_size=84) -> None:
    """P2 三杀手数字芯片(P9 大数字也复用)。"""
    add_text(slide, x_in, y_in, w_in, h_in * 0.55, num,
             font=FONT_MONO, size=num_size, color=color, bold=True)
    add_text(slide, x_in, y_in + h_in * 0.55, w_in, h_in * 0.2, unit,
             font=FONT_MONO, size=20, color=color)
    add_text(slide, x_in, y_in + h_in * 0.75, w_in, h_in * 0.2, label,
             font=FONT_CN, size=14, color=TEXT_SECONDARY)


def stat_card(slide, x_in, y_in, w_in, h_in, num, label, *,
              color=DANGER, num_size=56, label_size=14,
              fill=BG_PANEL, border_pt=2) -> None:
    """P2 警示/反衬数字卡:大号 mono 数字 + 中文副标签,圆角矩形边框。

    DANGER 用于警示场景,SUCCESS 用于反衬/正向场景。
    """
    add_rect(slide, x_in, y_in, w_in, h_in,
             fill=fill, line_color=color, line_width_pt=border_pt,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x_in, y_in + 0.1, w_in, h_in * 0.5, num,
             font=FONT_MONO, size=num_size, color=color, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, x_in, y_in + h_in * 0.6, w_in, h_in * 0.3, label,
             font=FONT_CN, size=label_size, color=TEXT_SECONDARY,
             align=PP_ALIGN.CENTER)


def chip_text(slide, x_in, y_in, w_in, h_in, text, *,
              color=SUCCESS, fill=BG_PANEL, line_width_pt=1.5,
              font=FONT_MONO, size=11, bold=True) -> None:
    """小徽章 / 标签:圆角矩形 + 单行 mono 字。

    用于:✓ 标记 / 章节金标 / 演示场景标注 / 路线图节点标签。
    """
    add_rect(slide, x_in, y_in, w_in, h_in,
             fill=fill, line_color=color, line_width_pt=line_width_pt,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x_in, y_in, w_in, h_in, text,
             font=font, size=size, color=color, bold=bold,
             align=PP_ALIGN.CENTER)


def feature_card(slide, x_in, y_in, w_in, h_in, icon, name, sub, *,
                 color=PRIMARY, fill=BG_PANEL,
                 icon_size=18, name_size=15, sub_size=11,
                 border_pt=2) -> None:
    """功能卡(图标 + 标题 + 副标题,彩色边框)。通用架构/功能网格单元。

    用于:P3 三个组件、P4 6 大功能、P5 4 starter、P9 3 仓库 等。
    """
    add_rect(slide, x_in, y_in, w_in, h_in,
             fill=fill, line_color=color, line_width_pt=border_pt,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x_in, y_in + 0.05, w_in, 0.4, icon,
             font=FONT_MONO, size=icon_size, color=color,
             align=PP_ALIGN.LEFT)
    add_text(slide, x_in, y_in + 0.35, w_in, 0.5, name,
             font=FONT_CN, size=name_size, color=color, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(slide, x_in, y_in + 0.85, w_in, h_in - 0.95, sub,
             font=FONT_CN, size=sub_size, color=TEXT_SECONDARY,
             align=PP_ALIGN.LEFT)


def bullet_list(slide, x_in, y_in, w_in, line_h_in, items, *,
                marker="✓", color=SUCCESS, size=14,
                text_color=None) -> None:
    """竖向 ✓ 列表 / → 列表 / 数字列表。

    items: 字符串列表(已含 marker 字符,或留空由本函数统一加)。
    line_h_in: 行高(英寸)
    text_color: 默认走 TEXT_PRIMARY
    """
    tc = text_color or TEXT_PRIMARY
    for i, line in enumerate(items):
        if marker and not line.startswith(marker) and not line.startswith(("→", "①", "②", "③", "④", "⑤", "⑥", "⑦", "⑧")):
            text = f"{marker} {line}"
        else:
            text = line
        add_text(slide, x_in, y_in + i * line_h_in,
                 w_in, line_h_in, text,
                 font=FONT_MONO, size=size, color=color, bold=True,
                 align=PP_ALIGN.LEFT)
        # 第二行后续用 tc 渲染
        # 注:为了简单,主色统一用 color;若需要区分,使用各自 add_text


def with_alpha(shape, alpha_pct) -> None:
    """给已添加的 shape 套 alpha(0-100),通过 a:solidFill/a:alpha 实现。

    用于:iso_lego highlight、装饰水印、半透明覆盖层。
    """
    solidFill = shape._element.spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    # 删旧的 alpha
    for old in solidFill.findall(qn("a:alpha")):
        solidFill.remove(old)
    a = etree.SubElement(solidFill, qn("a:alpha"))
    a.set("val", str(int(alpha_pct * 1000)))


def arrow_line(slide, x1, y1, x2, y2, *, color=TEXT_SECONDARY, weight_pt=1.0):
    """画一条直线(架构图连线用)。"""
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight_pt)
    return ln


def iso_lego(slide, x_in, y_in, w_in, h_in, color, *, studs=True, highlight=False):
    """画一个 2.5D 乐高积木(圆角矩形主体 + 顶部 4 圆钉)。

    适用于 v1/v2 风格。返回主形状。
    """
    main = add_rect(slide, x_in, y_in, w_in, h_in,
                     fill=color, line_color=None,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if studs:
        stud_r = min(w_in, h_in) * STUD_RATIO
        stud_y = y_in - stud_r * STUD_Y_OFFSET_RATIO
        stud_xs = [x_in + w_in * r for r in STUD_X_POSITIONS]
        for sx in stud_xs:
            stud = add_rect(slide, sx - stud_r / 2, stud_y,
                             stud_r, stud_r, fill=color,
                             shape=MSO_SHAPE.OVAL)
    if highlight:
        hl = add_rect(slide, x_in, y_in, w_in, h_in * 0.12,
                       fill=WHITE)
        # 30% alpha
        solidFill = hl._element.spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            alpha = etree.SubElement(solidFill, qn("a:alpha"))
            alpha.set("val", "30000")
    return main


# ====================================================================
# E. 文件 I/O
# ====================================================================

def replace_pptx_media(src: Path, dst: Path, target: str, data: bytes) -> None:
    """解压 pptx 替换指定 media 文件后重新打包。

    target 形如 'ppt/media/image2.png'。
    若 src == dst,先写到临时文件再 rename(避免同读同写错误)。
    """
    src = Path(src)
    dst = Path(dst)
    if src == dst:
        # 同读同写:用临时文件
        tmp = src.with_suffix(src.suffix + ".tmp")
        with zipfile.ZipFile(src, "r") as zin:
            with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    buf = zin.read(item.filename)
                    if item.filename == target:
                        buf = data
                    zout.writestr(item, buf)
        # 原子替换
        if dst.exists():
            dst.unlink()
        tmp.rename(dst)
    else:
        if dst.exists():
            dst.unlink()
        with zipfile.ZipFile(src, "r") as zin:
            with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    buf = zin.read(item.filename)
                    if item.filename == target:
                        buf = data
                    zout.writestr(item, buf)


def find_dominant_bg(image_path) -> tuple[int, int, int]:
    """从图片文件找全图最常见的颜色(主背景色)。"""
    from PIL import Image
    import numpy as np
    img = np.array(Image.open(image_path).convert("RGB"))
    sample = img[::20, ::20].reshape(-1, 3)
    quantized = (sample // 16) * 16
    c = Counter(map(tuple, quantized))
    top = c.most_common(1)[0][0]
    return (int(top[0]), int(top[1]), int(top[2]))


def cover_region(img, region, bg_rgb, pad=8):
    """物理覆盖图片指定矩形区域为背景色,边缘用 8px 高斯模糊软化。

    img: PIL.Image (RGB 或 RGBA)
    region: (x1, y1, x2, y2) 像素坐标
    bg_rgb: (r, g, b) 覆盖色
    返回: 新 Image (RGBA)
    """
    from PIL import Image, ImageDraw, ImageFilter
    x1, y1, x2, y2 = region
    W, H = img.size
    x1, y1 = max(0, x1 - pad), max(0, y1 - pad)
    x2, y2 = min(W, x2 + pad), min(H, y2 + pad)

    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    draw.rectangle([(x1, y1), (x2, y2)], fill=(*bg_rgb, 255))
    overlay = overlay.filter(ImageFilter.GaussianBlur(radius=8))
    return Image.alpha_composite(img.convert("RGBA"), overlay)


def qa_image_check(path) -> tuple[bool, list[str]]:
    """图片质量门 A:尺寸 / 文件大小 / 标准差 / 长宽比。

    阈值(可在 tools 顶部常量调整):
        - 短边 ≥ QA_MIN_SHORT_SIDE (720px)
        - 文件 ≥ QA_MIN_FILE_BYTES, ≤ QA_MAX_FILE_BYTES
        - 像素值标准差 ≥ QA_MIN_STD (避免纯色 / 全白 / 全黑)
        - 宽高比 ∈ [QA_MIN_RATIO, QA_MAX_RATIO] (避免极窄 / 极扁)

    Returns:
        (passed: bool, reasons: list[str])
    """
    from PIL import Image
    import numpy as np

    p = Path(path)
    reasons: list[str] = []

    if not p.exists():
        return False, [f"文件不存在: {p}"]

    try:
        img = Image.open(p).convert("RGB")
    except Exception as e:
        return False, [f"无法打开: {e}"]

    w, h = img.size
    short_side = min(w, h)
    if short_side < QA_MIN_SHORT_SIDE:
        reasons.append(f"短边过小: {short_side} < {QA_MIN_SHORT_SIDE}")

    size = p.stat().st_size
    if size < QA_MIN_FILE_BYTES:
        reasons.append(f"文件过小: {size//1024}KB < {QA_MIN_FILE_BYTES//1024}KB")
    if size > QA_MAX_FILE_BYTES:
        reasons.append(f"文件过大: {size//1024}KB > {QA_MAX_FILE_BYTES//1024}KB")

    arr = np.array(img)
    std = float(arr.std())
    if std < QA_MIN_STD:
        reasons.append(f"图像过平(std={std:.1f}<{QA_MIN_STD}),疑似纯色/全黑")

    ratio = w / max(h, 1)
    if ratio < QA_MIN_RATIO or ratio > QA_MAX_RATIO:
        reasons.append(f"宽高比异常: {ratio:.2f} (允许 {QA_MIN_RATIO}~{QA_MAX_RATIO})")

    return len(reasons) == 0, reasons


def new_presentation() -> Presentation:
    """创建一个 16:9 空 Presentation。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def find_media_target(prs, hint="image", ext="png") -> str | None:
    """在 prs 全部 slide 中查找匹配的 media 路径。

    hint: 'image' / 'image1' / 'logo' 等子串
    ext: 后缀(默认 png)

    Returns:
        'ppt/media/image3.png' 这种完整路径,找不到返回 None。
    """
    for slide in prs.slides:
        slide_part = slide.part
        for shape in slide.shapes:
            blip = shape._element.find(".//" + qn("a:blip"))
            if blip is None:
                continue
            rid = blip.get(qn("r:embed"))
            if not rid:
                continue
            try:
                partname = str(slide_part.related_part(rid).partname)
            except KeyError:
                continue
            if "media/" in partname and hint in partname and partname.endswith(f".{ext}"):
                return partname
    return None


def qa_pptx_images(pptx_path) -> tuple[bool, list[tuple[int, str, list[str]]]]:
    """把 pptx 内嵌的所有图片(去重)导出到临时目录,跑 qa_image_check。

    Returns:
        (passed, [(idx, filename, reasons), ...])
    """
    import tempfile
    import shutil
    import zipfile

    p = Path(pptx_path)
    if not p.exists():
        return False, [(0, str(p), [f"文件不存在: {p}"])]

    tmpdir = Path(tempfile.mkdtemp(prefix="qa_pptx_img_"))
    try:
        seen: set[str] = set()
        with zipfile.ZipFile(p) as zf:
            for name in zf.namelist():
                if name.startswith("ppt/media/") and name.endswith((".png", ".jpg", ".jpeg")):
                    base = name.split("/")[-1]
                    if base in seen:
                        continue
                    seen.add(base)
                    data = zf.read(name)
                    (tmpdir / base).write_bytes(data)

        failures: list[tuple[int, str, list[str]]] = []
        for idx, fn in enumerate(sorted(seen), 1):
            fpath = tmpdir / fn
            ok, reasons = qa_image_check(fpath)
            if not ok:
                failures.append((idx, fn, reasons))

        return len(failures) == 0, failures
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


def preview_grid(
    images_dir: Path,
    out_path: Path,
    *,
    pattern: str = "page-*.png",
    cols: int = 2,
    gap: int = 20,
    label_h: int = 40,
    bg: tuple[int, int, int] = (32, 32, 32),
    label_format: str = "P{idx:02d}",
) -> Path:
    """把目录下匹配 pattern 的图片拼成 cols×N 网格,用于整体预览风格。

    用法:
        preview_grid(
            Path("docs/ppt/images"),
            Path("docs/ppt/preview/_grid.png"),
        )
        # 自动找 page-01-cover.png / page-02-pain.png / ... 拼 2x5 网格

    Returns:
        out_path(已写入磁盘)
    """
    from PIL import Image, ImageDraw, ImageFont

    images_dir = Path(images_dir)
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    paths = sorted(images_dir.glob(pattern))
    if not paths:
        raise FileNotFoundError(
            f"在 {images_dir} 找不到匹配 {pattern!r} 的图片"
        )

    imgs = [Image.open(p) for p in paths]
    w, h = imgs[0].size
    # 如果图片尺寸不一致,统一缩放到最小
    sizes = {im.size for im in imgs}
    if len(sizes) > 1:
        imgs = [im.resize((w, h)) for im in imgs]

    rows = (len(imgs) + cols - 1) // cols
    W = cols * w + (cols + 1) * gap
    H = rows * (h + label_h) + (rows + 1) * gap

    canvas = Image.new("RGB", (W, H), bg)
    draw = ImageDraw.Draw(canvas)

    # 中文字体
    try:
        font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 24)
    except Exception:
        try:
            font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 24)
        except Exception:
            font = ImageFont.load_default()

    for i, im in enumerate(imgs, 1):
        col = (i - 1) % cols
        row = (i - 1) // cols
        x = gap + col * (w + gap)
        y = gap + row * (h + label_h + gap)
        draw.text((x, y), label_format.format(idx=i), fill="white", font=font)
        canvas.paste(im, (x, y + label_h))

    canvas.save(str(out_path))
    return out_path


# ====================================================================
# E. 万相 API(阿里云百炼 DashScope 2.7)
# ====================================================================

# 万相 2.7 同步生成端点(北京地域)
WANX_GENERATE_URL = (
    "https://dashscope.aliyuncs.com/api/v1/services/aigc/"
    "multimodal-generation/generation"
)


def _wanx_api_key() -> str:
    """从环境变量取 API key;未设置则抛错。"""
    k = os.environ.get("DASHSCOPE_API_KEY")
    if not k:
        raise RuntimeError(
            "DASHSCOPE_API_KEY not set. "
            "Run: export DASHSCOPE_API_KEY=sk-xxx"
        )
    return k


def _wanx_post(url: str, data: dict, timeout: int = 120) -> dict:
    """统一 HTTP POST,返回 dict;失败抛错并打印详情。"""
    headers = {
        "Authorization": f"Bearer {_wanx_api_key()}",
        "Content-Type": "application/json",
    }
    body = json.dumps(data).encode("utf-8")
    req = urllib.request.Request(url, data=body, headers=headers, method="POST")
    try:
        with urllib.request.urlopen(req, timeout=timeout) as r:
            return json.loads(r.read())
    except urllib.error.HTTPError as e:
        body_text = e.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"wanx HTTP {e.code}\n{body_text}") from e
    except urllib.error.URLError as e:
        raise RuntimeError(f"wanx network error: {e.reason}") from e


def wanx_generate(
    prompt: str,
    *,
    model: str = "wan2.7-image",
    size: str = "1280*720",
    n: int = 1,
    watermark: bool = False,
) -> str:
    """万相 2.7 同步生成图片,返回 image_url(24h 有效)。

    用法:
        url = wanx_generate("a cat sitting on a chair")
        # url = "https://dashscope-.../xxx.png"
    """
    payload = {
        "model": model,
        "input": {
            "messages": [
                {"role": "user", "content": [{"text": prompt}]}
            ]
        },
        "parameters": {"size": size, "n": n, "watermark": watermark},
    }
    rsp = _wanx_post(WANX_GENERATE_URL, payload)
    try:
        return rsp["output"]["choices"][0]["message"]["content"][0]["image"]
    except (KeyError, IndexError, TypeError) as e:
        raise RuntimeError(
            f"wanx_generate: unexpected response: "
            f"{json.dumps(rsp, ensure_ascii=False, indent=2)}"
        ) from e


def wanx_upload(local_path: Path) -> str:
    """上传本地图片到 DashScope Files,返回带签名 OSS URL。

    image_edit / imageedit 等需要 file URL(非 file_id)。
    """
    from dashscope import Files

    up = Files.upload(file_path=str(local_path), purpose="image_edit")
    if up.status_code != 200:
        raise RuntimeError(f"upload failed: {up.status_code} {up.message}")
    file_id = up.output["uploaded_files"][0]["file_id"]
    g = Files.get(file_id=file_id)
    if g.status_code != 200:
        raise RuntimeError(f"Files.get failed: {g.message}")
    url = g.output.get("url")
    if not url:
        raise RuntimeError(f"no url in get response: {g.output}")
    return url


def wanx_edit(
    base_image_url: str,
    prompt: str,
    *,
    mask_image_url: str | None = None,
    function: str = "description_edit",
    model: str = "wanx2.1-imageedit",
    n: int = 1,
) -> str:
    """万相图像编辑(2.1 imageedit),返回 image_url。

    function 取值:
        - description_edit             全图重绘
        - description_edit_with_mask   局部重绘(需 mask_image_url)
        - stylization_all / stylization_local
        - remove_watermark
        - expand / super_resolution / colorization
    """
    from dashscope import ImageSynthesis

    kwargs = dict(
        model=model,
        function=function,
        base_image_url=base_image_url,
        prompt=prompt,
        n=n,
    )
    if mask_image_url is not None:
        kwargs["mask_image_url"] = mask_image_url
    rsp = ImageSynthesis.call(**kwargs)
    if rsp.status_code != 200:
        raise RuntimeError(f"wanx_edit failed: {rsp.status_code} {rsp.message}")
    out_url = rsp.output["results"][0].get("url")
    if not out_url:
        raise RuntimeError(f"no url in result: {rsp.output}")
    return out_url


def wanx_download(url: str, dest: Path, timeout: int = 60) -> Path:
    """下载图片 URL 到本地路径,返回 dest。"""
    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    req = urllib.request.Request(url, headers={"User-Agent": "wanx-download/1.0"})
    with urllib.request.urlopen(req, timeout=timeout) as r, open(dest, "wb") as f:
        while True:
            chunk = r.read(8192)
            if not chunk:
                break
            f.write(chunk)
    return dest
