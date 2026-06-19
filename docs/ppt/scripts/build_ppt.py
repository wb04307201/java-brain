#!/usr/bin/env python3
"""根据 OUTLINE_v1.md 组装 PPT。

使用 docs/ppt/images/ 下的图片创建 16:9 PPT。
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx.util import Inches
from tools import (
    new_presentation,
    apply_chrome,
)


# 页面配置：(页码, 主题, 章节标题)
PAGES = [
    (1, "page-01-cover.png", "JavaBrain"),
    (2, "page-02-pain.png", "企业 AI 落地的 4 道坎"),
    (3, "page-03-position.png", "定位 + 组成"),
    (4, "page-04-loom.png", "灵梭 · 6 大功能模块"),
    (5, "page-05-sql-forge.png", "SQL 工坊 · 4 starter + 6 大功能"),
    (6, None, "录屏 1：一句话出分析报告"),
    (7, None, "录屏 2：一句话生成 CRUD"),
    (8, "page-08-comparison.png", "差异化对比"),
    (9, "page-09-roadmap.png", "开源 + 路线图"),
    (10, "page-10-ending.png", "结尾"),
]


def main():
    """主流程。"""
    # 设置路径
    script_dir = Path(__file__).parent  # docs/ppt/scripts/
    ppt_dir = script_dir.parent  # docs/ppt/
    images_dir = ppt_dir / "images"
    output_dir = ppt_dir

    print("=" * 70)
    print("组装 PPT")
    print("=" * 70)

    # 创建 Presentation
    prs = new_presentation()

    total = len(PAGES)

    for page_num, image_name, section in PAGES:
        print(f"\n添加第 {page_num}/{total} 页: {section}")

        # 添加空白 slide
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 应用 chrome（背景 + 网格 + 角标 + 章节标）
        apply_chrome(slide, page_num, total, section)

        # 如果有图片，添加到 slide
        if image_name:
            image_path = images_dir / image_name
            if image_path.exists():
                # 添加图片到 slide 中央
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(1.5),  # x
                    Inches(1.0),  # y
                    Inches(10.333),  # w
                    Inches(5.5),  # h
                )
                print(f"  [OK] 添加图片: {image_name}")
            else:
                print(f"  [WARN] 图片不存在: {image_path}")

    # 保存
    output_path = output_dir / "javabrain_v2.pptx"
    prs.save(str(output_path))

    print("\n" + "=" * 70)
    print(f"[OK] PPT 已保存: {output_path}")
    print(f"共 {total} 页")
    print("=" * 70)

    return 0


if __name__ == "__main__":
    sys.exit(main())
