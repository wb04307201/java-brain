#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""JavaBrain 答辩 PPT 生成器 - v2(白底 2.5D 乐高)"""
from __future__ import annotations

import sys
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ============== 路径 ==============
SCRIPT_DIR = Path(__file__).resolve().parent
PPT_DIR = SCRIPT_DIR.parent
OUTPUT = PPT_DIR / "javabrain.pptx"

# ============== 幻灯片尺寸(16:9) ==============
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============== 配色(7 色 + 4 中性,见 spec §4.1) ==============
JAVA_BLUE = RGBColor(0x00, 0x73, 0x96)
AI_PURPLE = RGBColor(0x6B, 0x46, 0xC1)
SEMANTIC_RED = RGBColor(0xDC, 0x26, 0x26)
SEMANTIC_GREEN = RGBColor(0x10, 0xB9, 0x81)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
TEXT_PRIMARY = RGBColor(0x1F, 0x29, 0x37)
TEXT_SECONDARY = RGBColor(0x6B, 0x72, 0x80)
BG_LIGHT = RGBColor(0xFA, 0xFB, 0xFC)
DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)
GOLD_BG = RGBColor(0xFE, 0xF3, 0xC7)
GREEN_BG = RGBColor(0xD1, 0xFA, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ============== 字体 ==============
FONT_CN = "阿里巴巴普惠体 3.0"
FONT_EN = "Inter"
FONT_MONO = "JetBrains Mono"

# ============== 动画类型(见 spec §5.2) ==============
ANIM_PRESETS = {
    "fade_in":   ("entr", "10"),
    "fade_out":  ("exit", "11"),
    "zoom_in":   ("entr", "23"),
    "appear":    ("entr", "1"),
    "fly_left":  ("entr", "3"),
    "fly_right": ("entr", "2"),
    "fly_top":   ("entr", "5"),
    "fly_bot":   ("entr", "4"),
    "pulse":     ("emph", "26"),
    "spin":      ("emph", "8"),
}


def add_anim(slide, shape, anim_type, *, delay_ms=0, dur_ms=500, loop=False):
    """注入 PowerPoint 动画节点(累积到 slide 唯一 <p:timing>)。

    关键:PowerPoint / OnlyOffice 要求每张 slide 只能有 1 个 <p:timing> 元素,
    所有动画必须嵌套在同一个 timing 的 tnLst 里。多次调用 add_anim 时
    找到/创建 slide 的 <p:timing>,把新动画 par 追加到 tnLst 的内层 par 里。

    anim_type: 见 ANIM_PRESETS 的 key
    loop=True: 生成 repeatCount="indefinite"(持续到翻页,转视频核心)
    fill="hold" 保持动画间形状可见;loop=True 时必需
    """
    if anim_type not in ANIM_PRESETS:
        raise ValueError(
            f"unknown anim_type {anim_type!r}; valid: {list(ANIM_PRESETS)}"
        )
    preset_class, preset_id = ANIM_PRESETS[anim_type]
    shape_id = shape.shape_id
    repeat_attr = ' repeatCount="indefinite"' if loop else ""

    # 找 slide 已有的 <p:timing>(用 ns0 命名空间别名)
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    timing = slide._element.find("p:timing", nsmap)
    if timing is None:
        # 第一次调用:创建 <p:timing> + tnLst + 顶层 par(顺序播放)
        timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:pPr/>
    </p:par>
  </p:tnLst>
</p:timing>'''
        timing = etree.fromstring(timing_xml)
        slide._element.append(timing)
    # 顶层 par(顺序容器,所有动画在此顺次排)
    top_par = timing.find("p:tnLst/p:par", nsmap)
    # 在 top_par 里追加一个并列 par(每个动画一个并列 par → 全部同时开始,各自带不同 begin)
    anim_par_xml = f'''<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cond><p:tnCond val="0"/></p:cond>
  <p:par>
    <p:cTn id="{shape_id + 1}" nodeType="withEffect" dur="{dur_ms}" begin="{delay_ms}ms" fill="hold"{repeat_attr}>
      <p:animEffect presetClass="{preset_class}" presetId="{preset_id}" dur="{dur_ms}"/>
    </p:cTn>
  </p:par>
</p:par>'''
    anim_par = etree.fromstring(anim_par_xml)
    top_par.append(anim_par)


def iso_lego(slide, x_in, y_in, w_in, h_in, color, *, studs=True, highlight=False):
    """画一个 2.5D 乐高积木(圆角矩形主体 + 顶部 4 圆钉)。
    返回主形状(用于 add_anim)。
    本版本第一版先简化:不做高光条和软阴影(后续可加)。
    """
    In = Inches
    main = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   In(x_in), In(y_in), In(w_in), In(h_in))
    main.fill.solid()
    main.fill.fore_color.rgb = color
    main.line.fill.background()
    if studs:
        stud_r = min(w_in, h_in) * 0.16
        stud_y = y_in - stud_r * 0.5
        stud_xs = [x_in + w_in * 0.18, x_in + w_in * 0.40,
                   x_in + w_in * 0.60, x_in + w_in * 0.82]
        for sx in stud_xs:
            stud = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                           In(sx - stud_r / 2), In(stud_y),
                                           In(stud_r), In(stud_r))
            stud.fill.solid()
            stud.fill.fore_color.rgb = color
            stud.line.fill.background()
    if highlight:
        hl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     In(x_in), In(y_in), In(w_in), In(h_in * 0.12))
        hl.fill.solid()
        hl.fill.fore_color.rgb = WHITE
        hl.line.fill.background()
        # 设置 30% alpha(spec §4.4 "白色 alpha 30%")
        from lxml import etree as _etree
        solidFill = hl._element.spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            alpha = _etree.SubElement(solidFill, qn("a:alpha"))
            alpha.set("val", "30000")  # 30000 = 30%
    return main


def styled_text(slide, x, y, w, h, text, *,
                font=None, size=18, color=None, bold=False, center=True):
    """创建并配置一个文本框。返回 textbox shape。
    font/color 默认:font=FONT_CN, color=TEXT_PRIMARY
    center=True → PP_ALIGN.CENTER
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True  # 避免自动收缩字号
    tf.text = text
    p = tf.paragraphs[0]
    if center:
        p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = font if font else FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color else TEXT_PRIMARY
    return tb


def _add_picture_with_alpha(slide, image_path, x, y, w, h, alpha_pct):
    """Add picture with overall alpha (0-100). alpha_pct=100 fully opaque, smaller = more transparent.
    python-pptx has no built-in API; we manipulate the blip's alphaModFix via XML.
    """
    from PIL import Image as _PILImage
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(w), Inches(h))
    # 设置图片整体透明度
    sppr = pic._element.spPr
    # 添加 alphaModFix 到 blip
    blip = pic._element.blipFill.blip
    # 在 blip 元素上设置 alpha(以 1000 为基数,如 30% = 30000)
    alpha_val = int(alpha_pct * 1000)
    # 找到 blip 元素的 a:alphaModFix 子元素,如果没有就加
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    alpha = blip.find("a:alphaModFix", nsmap)
    if alpha is None:
        alpha = etree.SubElement(blip, qn("a:alphaModFix"))
    alpha.set("amt", str(alpha_val))
    return pic


def add_page_number(slide, current, total=12, color=TEXT_SECONDARY, size=10):
    """在右下角加 'XX/12' 页码(灰色小字,品牌系统一致)。"""
    styled_text(slide, 11.6, 7.05, 1.4, 0.3,
                f"{current:02d} / {total}",
                font=FONT_EN, size=size, color=color)


def slide_1_cover(prs):
    """P1 封面:AI 神经网络背景图(半透明)+ 3 圈光晕节点 + JavaBrain 大字 + 副标 + 金钩。
    替换原"3 乐高"封面,5 fade_in。
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # === 1) JavaBrain 大脑配图(顶部 3.3 英寸,35% alpha 作为背景) ===
    # 图源: docs/ppt/images/ai/p1v2-d1-brain-3streams.png (紫发光大脑 + 蓝/绿/紫 3 股能量流)
    img_path = PPT_DIR / "images" / "ai" / "p1v2-d1-brain-3streams.png"
    pic = _add_picture_with_alpha(s, img_path, 3.0, 0.1, 7.5, 3.3, alpha_pct=38)
    # 派发到最底层(zorder),作为背景
    spTree = s.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)  # nvGrpSpPr 之后第 1 个位置

    # === 2) JavaBrain 大字(72pt 居中)===
    tb_title = styled_text(s, 1.0, 3.5, 11.333, 1.4,
                           "JavaBrain",
                           font=FONT_EN, size=72, color=TEXT_PRIMARY, bold=True)

    # === 3) 副标 ===
    tb_sub = styled_text(s, 1.0, 4.9, 11.333, 0.5,
                          "让您的系统瞬间拥有“思考”与“执行”的智能大脑",
                          font=FONT_CN, size=24, color=TEXT_SECONDARY)

    # === 4) 3 组件标签(横向,Java 蓝 / AI 紫 / 绿)===
    sat_y_lbl = 5.6
    sat_xs = [3.5, 6.167, 8.833]
    sat_colors = [JAVA_BLUE, AI_PURPLE, SEMANTIC_GREEN]
    sat_shapes = []
    for sx, sc in zip(sat_xs, sat_colors):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(sx), Inches(sat_y_lbl),
                                  Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = sc
        dot.line.fill.background()
        sat_shapes.append(dot)
    labels = [("灵梭", JAVA_BLUE), ("SQL工坊", AI_PURPLE), ("SQL工坊 MCP", SEMANTIC_GREEN)]
    for sx, (txt, col) in zip(sat_xs, labels):
        tb = styled_text(s, sx + 0.25, sat_y_lbl - 0.07, 2.0, 0.4,
                          txt, font=FONT_CN, size=14, color=col, bold=True)
        sat_shapes.append(tb)

    # === 5) 金钩"3 分钟接入 AI" — 数字钩,呼应 P2 时间轴 ===
    hook_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4.0), Inches(6.4),
                                   Inches(5.333), Inches(0.6))
    hook_box.fill.solid()
    hook_box.fill.fore_color.rgb = GOLD
    hook_box.line.fill.background()
    styled_text(s, 4.0, 6.45, 5.333, 0.5,
                "3 分钟接入 AI",
                font=FONT_CN, size=22, color=WHITE, bold=True)

    # === 6) 6 动画(大脑 zoom_in 缩放入场 + 标题副标 + 3 组件 + 金钩 zoom_in 模拟打字机)===
    add_anim(s, pic._element,  "zoom_in", delay_ms=0,    dur_ms=800)  # 大脑 0.9→1.0
    add_anim(s, tb_title,      "fade_in", delay_ms=600,  dur_ms=500)
    add_anim(s, tb_sub,        "fade_in", delay_ms=1000, dur_ms=500)
    add_anim(s, sat_shapes[0], "fade_in", delay_ms=1400, dur_ms=500)
    add_anim(s, sat_shapes[3], "fade_in", delay_ms=1600, dur_ms=500)
    add_anim(s, hook_box,      "zoom_in", delay_ms=2000, dur_ms=600)  # 金钩 0.5x→1.0x 弹出


def slide_2_pain(prs):
    """P2 痛点:横向时间轴 + 6 槽位(3 红→3 绿)+ 6 张配图 + 中心"1440×"金钩 pulse。
    替换原"3 红 vs 3 绿平铺"布局。
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # === 标题 ===
    tb_t = styled_text(s, 0.5, 0.3, 12.333, 0.6,
                       "接入 AI 的 3 个痛点  vs  JavaBrain 的 3 个解法",
                       size=24, bold=True)

    # === 时间轴线(从左红→中金→右绿,横向贯穿中部)===
    # 整条 12 in 长 × 0.06 in 厚
    axis_y = 4.0
    axis_x0 = 0.667
    axis_w = 12.0
    # 主体白底轴线
    axis = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(axis_x0), Inches(axis_y),
                              Inches(axis_w), Inches(0.06))
    axis.fill.solid()
    axis.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)  # 浅灰
    axis.line.fill.background()

    # 红段(左 1/2)
    red_seg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(axis_x0), Inches(axis_y),
                                  Inches(6.0), Inches(0.06))
    red_seg.fill.solid()
    red_seg.fill.fore_color.rgb = SEMANTIC_RED
    red_seg.line.fill.background()

    # 金段(中间 1 in,锚点附近)
    gold_seg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(axis_x0 + 5.85), Inches(axis_y),
                                   Inches(1.5), Inches(0.06))
    gold_seg.fill.solid()
    gold_seg.fill.fore_color.rgb = GOLD
    gold_seg.line.fill.background()

    # 绿段(右 1/2)
    green_seg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(axis_x0 + 6.0), Inches(axis_y),
                                    Inches(6.0), Inches(0.06))
    green_seg.fill.solid()
    green_seg.fill.fore_color.rgb = SEMANTIC_GREEN
    green_seg.line.fill.background()

    # === 6 槽位(3 红 + 3 绿)===
    # 槽位布局:6 个圆形锚点沿时间轴,X 间距 = 12/5 = 2.4
    slot_xs = [axis_x0 + i * 2.4 for i in range(6)]
    # 槽位配置:(img_filename, 数字, 文字, 颜色)
    # 6 张图已用 PIL 裁掉文字/边框/按钮
    slots = [
        # 红段(传统)— 2.5D 几何图标
        ("p2-red-1-hourglass.png",      "3 月",   "集成 AI",  SEMANTIC_RED),
        ("p2-red-2-bifurcation.png",    "5 天",   "训练数据",  SEMANTIC_RED),
        ("p2-red-3-tangled.png",        "3 天",   "上线联调",  SEMANTIC_RED),
        # 绿段(JavaBrain)
        ("p2-green-1-lightning.png",    "3 分",   "一键集成",  SEMANTIC_GREEN),
        ("p2-green-2-stopwatch.png",    "90 秒",  "出报告",    SEMANTIC_GREEN),
        ("p2-green-3-puzzle.png",       "10 分",  "出页面",    SEMANTIC_GREEN),
    ]

    anim_targets = []  # 用于动画的对象
    for i, (img, big, sub, col) in enumerate(slots):
        cx = slot_xs[i]
        # 1) 上方配图(1.2 in × 1.2 in)
        img_w = 1.2
        img_path = PPT_DIR / "images" / "ai" / img
        if img_path.exists():
            pic = s.shapes.add_picture(str(img_path),
                                       Inches(cx - img_w/2), Inches(1.4),
                                       Inches(img_w), Inches(img_w))
            anim_targets.append(pic)
        # 2) 锚点(时间轴上的圆点)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(cx - 0.12), Inches(axis_y - 0.09),
                                  Inches(0.24), Inches(0.24))
        dot.fill.solid()
        dot.fill.fore_color.rgb = col
        dot.line.color.rgb = WHITE
        dot.line.width = Pt(2)
        # 3) 大数字(锚点下方)
        tb_big = styled_text(s, cx - 0.7, axis_y + 0.3, 1.4, 0.6,
                              big, font=FONT_EN, size=24, color=col, bold=True)
        # 4) 小字说明(数字下方)
        tb_sub2 = styled_text(s, cx - 0.8, axis_y + 0.95, 1.6, 0.4,
                               sub, font=FONT_CN, size=14, color=TEXT_SECONDARY)
        anim_targets.extend([dot, tb_big, tb_sub2])

    # === 中心金钩"1440×"(金圆 + 数字 + 副标)===
    # 位置:中心时间轴上(中点 = axis_x0 + 6.0)
    center_x = axis_x0 + 6.0
    # 大金圆
    big_circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(center_x - 0.55), Inches(axis_y - 0.55),
                                     Inches(1.1), Inches(1.1))
    big_circle.fill.solid()
    big_circle.fill.fore_color.rgb = GOLD_BG
    big_circle.line.color.rgb = GOLD
    big_circle.line.width = Pt(3)
    # 数字 1440×
    styled_text(s, center_x - 0.6, axis_y - 0.25, 1.2, 0.5,
                "1440×", font=FONT_EN, size=22, color=GOLD, bold=True)
    # 副标"效率提升"
    styled_text(s, center_x - 0.7, axis_y + 0.18, 1.4, 0.4,
                "效率提升", font=FONT_CN, size=11, color=GOLD, bold=True)

    # 杀手锏金句(底部)
    hook = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(3.5), Inches(6.3),
                              Inches(6.333), Inches(0.7))
    hook.fill.solid()
    hook.fill.fore_color.rgb = GOLD_BG
    hook.line.color.rgb = GOLD
    hook.line.width = Pt(3)
    styled_text(s, 3.5, 6.35, 6.333, 0.6,
                "★ 3 月 vs 3 分 — 从季度到分钟",
                size=20, color=GOLD, bold=True)

    # === 动画(策略 B:6 配图依次入场 + 1 数字串入场 + 1 中心金钩 pulse)===
    # 6 张图交错 350ms
    for i in range(min(6, len(anim_targets))):
        add_anim(s, anim_targets[i], "fade_in", delay_ms=i * 350, dur_ms=500)
    # 中心金钩 pulse(杀手锏持续脉冲)
    add_anim(s, hook, "pulse", delay_ms=3500, dur_ms=1500, loop=True)
    # 底部 6 文字组整体淡入(轻量,共 6 + 1 center + 1 hook = 8 入场 + 1 pulse)
    add_anim(s, big_circle, "fade_in", delay_ms=2200, dur_ms=500)


def slide_3_position(prs):
    """P3 定位:3 乐高(带标签)+ 4 行 ✓ + 首行金脉冲(策略 A:4 动画=3 入场 + 1 循环)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # 标题
    tb_t = styled_text(s, 0.5, 0.4, 12.333, 0.7,
                       "JavaBrain = 灵梭 + SQL工坊 + SQL工坊 MCP",
                       size=28, bold=True)

    # 依赖关系图(全宽,中部,100% alpha)— 3 核心 + 6 周边 + 连线
    # 高度压到 2.2 避免覆盖下面杀手锏(y=1.1-3.3)
    dep_img = PPT_DIR / "images" / "ai" / "p3_dependency.png"
    pic = _add_picture_with_alpha(s, dep_img, 3.5, 1.1, 6.333, 2.2, alpha_pct=100)
    spTree = s.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)

    # 3 组件标签(位置在图下方,y=3.3)
    styled_text(s, 3.5, 3.3, 2.0, 0.3, "灵梭 (AI 编排)",
                size=12, color=JAVA_BLUE, bold=True, center=True)
    styled_text(s, 5.7, 3.3, 2.0, 0.3, "SQL工坊 (数据底座)",
                size=12, color=AI_PURPLE, bold=True, center=True)
    styled_text(s, 7.8, 3.3, 2.0, 0.3, "SQL工坊 MCP (桥)",
                size=12, color=SEMANTIC_GREEN, bold=True, center=True)
    legos = [pic._element]  # 复用动画对象

    # 4 行杀手锏(4 个差异化优势,各对应 1 个组件的核心壁垒)
    # y=3.95 起点避免和上方组件标签重叠
    items = [
        ("★ 灵梭杀手锏:零配置接入 RAG / MCP / Skill · 50 行配置 → 1 行模板", GOLD, True),
        ("★ SQL工坊杀手锏:JSON CRUD + Calcite 联邦 + Amis 低代码 · 开发效率 +300%", TEXT_PRIMARY, False),
        ("★ SQL工坊 MCP 杀手锏:5 受限通道 · 数据不出企业 · 私有化部署", AI_PURPLE, False),
        ("★ 全部杀手锏:3 个依赖库可独立按需引入 · 不绑死 Spring Boot 生态", SEMANTIC_GREEN, False),
    ]
    item_boxes = []
    for i, (text, color, _) in enumerate(items):
        tb = styled_text(s, 0.7, 3.95 + i * 0.55, 11.9, 0.5,
                         text, size=16, color=color, bold=(i == 0))
        item_boxes.append(tb)

    # 4 动画:策略 A(标题入场 + 1 乐高入场 + 首行金行入场 + 1 循环)
    add_anim(s, tb_t, "fade_in", delay_ms=0, dur_ms=500)
    add_anim(s, legos[0], "fade_in", delay_ms=600, dur_ms=500)
    add_anim(s, item_boxes[0], "fade_in", delay_ms=1200, dur_ms=500)
    add_anim(s, item_boxes[0], "pulse", delay_ms=2000, dur_ms=1500, loop=True)


def slide_4a_loom_intro(prs):
    """P4-a 灵梭 · 3 步进阶(演进路线)+ AI 编排图作为视觉锚(alpha=100% 避免 OnlyOffice 渲染跳过)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # 标题
    styled_text(s, 0.5, 0.3, 12.333, 0.4,
                "灵梭 · LoomAgent",
                font=FONT_EN, size=16, color=JAVA_BLUE, bold=True)

    tb_t = styled_text(s, 0.5, 0.7, 12.333, 0.7,
                       "从 Spring AI 裸用到 JavaBrain 一体化",
                       size=28, bold=True)

    # 副标杀手锏
    styled_text(s, 0.5, 1.4, 12.333, 0.4,
                "★ 50 行配置 → 1 行模板,半年维护成本归零",
                size=15, color=GOLD, bold=True)

    # 3 步进阶横排(每步 = 大编号 + 标题 + 描述 + 时间标签)
    steps = [
        ("1", "裸 Spring AI", "写 50 行\nChatClient 配置", "5 天", DIVIDER, False),
        ("2", "灵梭 starter", "1 个 @Bean\n完成 RAG + MCP", "1 天", JAVA_BLUE, False),
        ("3", ".st 模板", "改 1 行字\nAI 行为变", "10 分钟", GOLD, True),
    ]
    step_cards = []
    sw, sh = 3.7, 3.0
    gap = 0.4
    sx0 = (13.333 - 3 * sw - 2 * gap) / 2
    sy0 = 2.0
    for i, (num, title, desc, time, accent, killer) in enumerate(steps):
        x = sx0 + i * (sw + gap)
        y = sy0
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(sw), Inches(sh))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = GOLD if killer else accent
        card.line.width = Pt(3) if killer else Pt(1.5)
        # 大编号(左上方,80pt)
        styled_text(s, x + 0.2, y + 0.1, 1.4, 1.4,
                    num, font=FONT_EN, size=72,
                    color=GOLD if killer else accent, bold=True)
        # 标题
        styled_text(s, x + 0.2, y + 1.5, sw - 0.4, 0.5,
                    title, size=20, color=accent, bold=True)
        # 描述
        styled_text(s, x + 0.2, y + 2.0, sw - 0.4, 0.6,
                    desc, size=13, color=TEXT_SECONDARY)
        # 时间标签(底部)
        tb_time = styled_text(s, x + 0.2, y + sh - 0.55, sw - 0.4, 0.4,
                    time, size=14, color=GOLD if killer else accent, bold=True)
        step_cards.append((card, tb_time))

    # 配图:左下 AI 编排图(放大 30% — 从 2.5x2.0 改 3.3x2.4,作为视觉锚)
    img_path = PPT_DIR / "images" / "ai" / "p4a-ai-orchestrate.png"
    pic = _add_picture_with_alpha(s, img_path, 0.2, 5.1, 3.3, 2.4, alpha_pct=100)
    spTree = s.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)

    # 8 动画:1 标题 + 3 步骤卡 + 1 .st 杀手锏 pulse + 3 时间标签
    add_anim(s, tb_t, "fade_in", delay_ms=0, dur_ms=500)
    for i, (card, tb_time) in enumerate(step_cards):
        add_anim(s, card, "fade_in", delay_ms=400 + i * 300, dur_ms=500)
        add_anim(s, tb_time, "fade_in", delay_ms=400 + i * 300, dur_ms=500)
    # .st 第 3 步杀手锏脉冲
    add_anim(s, step_cards[2][0], "pulse", delay_ms=2000, dur_ms=1500, loop=True)


def slide_4b_loom_modules(prs):
    """P4-b 灵梭 · 6 模块功能 + 2 chase pulse(策略 C:8 动画=6 入场 + 2 chase 循环)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # 标题
    styled_text(s, 0.5, 0.4, 12.333, 0.6,
                "灵梭 · 6 大功能模块",
                size=28, bold=True)

    # 6 卡片(2x3)— 杀手锏优先(Row1),基础能力(Row2)
    modules = [
        # Row 1:3 杀手锏(金边 + 金字)
        ("📚 RAG 知识库 ★", "Tika + JVector\n文档向量化检索", GOLD, True),
        ("🔧 MCP ★", "8+ 服务可热插拔\n工具热加载", GOLD, True),
        ("🧠 Skill ★", ".st 模板即提示词\n改 1 行 AI 行为变", GOLD, True),
        # Row 2:3 基础能力(浅灰边 + 蓝字)
        ("💬 对话 UI", "SSE + 思维链折叠\n流式输出", JAVA_BLUE, False),
        ("📁 文件管理", "磁盘 + H2 元数据\n会话持久化", JAVA_BLUE, False),
        ("🛠️ 内置工具", "时间 / Git / Maven\n沙箱执行", JAVA_BLUE, False),
    ]
    cards = []
    cw, ch = 4.0, 1.8
    gap_x, gap_y = 0.2, 0.2
    cx0 = (13.333 - 3 * cw - 2 * gap_x) / 2
    cy0 = 1.3
    for i, (name, desc, color, killer) in enumerate(modules):
        col, row = i % 3, i // 3
        x = cx0 + col * (cw + gap_x)
        y = cy0 + row * (ch + gap_y)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(cw), Inches(ch))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        if killer:
            card.line.color.rgb = GOLD
            card.line.width = Pt(3)
        else:
            card.line.color.rgb = DIVIDER
            card.line.width = Pt(1)
        # 名字
        styled_text(s, x, y + 0.2, cw, 0.6, name, size=20, color=color, bold=True)
        # 描述
        styled_text(s, x, y + 0.9, cw, 0.5, desc, size=14, color=TEXT_SECONDARY)
        cards.append(card)

    # 底部金句
    styled_text(s, 0.5, 6.5, 12.333, 0.5,
                "★ 改一个 .st 文件,AI 行为就变",
                size=20, color=GOLD, bold=True)

    # 8 动画:6 入场 + 2 chase(策略 C,chase 错开 2s)
    for i in range(6):
        add_anim(s, cards[i], "fade_in",
                 delay_ms=i * 200, dur_ms=500)
    # 2 杀手锏金边 chase(MCP + Skill)
    add_anim(s, cards[1], "pulse", delay_ms=10000, dur_ms=1500, loop=True)  # MCP
    add_anim(s, cards[2], "pulse", delay_ms=12000, dur_ms=1500, loop=True)  # Skill


def slide_5a_forge_intro(prs):
    """P5-a SQL工坊 · 3 步进阶(演进路线)+ Calcite 联邦图作为视觉锚(alpha=100%)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # 标题
    styled_text(s, 0.5, 0.3, 12.333, 0.4,
                "SQL工坊 · SQL Forge",
                font=FONT_EN, size=16, color=AI_PURPLE, bold=True)

    tb_t = styled_text(s, 0.5, 0.7, 12.333, 0.7,
                       "从 JDBC 散弹到 Calcite 联邦",
                       size=28, bold=True)

    # 副标杀手锏
    styled_text(s, 0.5, 1.4, 12.333, 0.4,
                "★ 1 套协议代替 5 套 ORM,数据库迁移零改代码",
                size=15, color=GOLD, bold=True)

    # 3 步进阶横排
    steps = [
        ("1", "裸 JDBC / MyBatis", "维护 N 套 CRUD\nN 套 ORM 配置", "30 天", DIVIDER, False),
        ("2", "SQL 工坊 starter", "1 套 JSON 协议\n跨库自动 JOIN", "1 天", AI_PURPLE, False),
        ("3", "sql-forge-mcp", "AI 5 受限通道\n数据不出企业", "0 直连", GOLD, True),
    ]
    step_cards = []
    sw, sh = 3.7, 3.0
    gap = 0.4
    sx0 = (13.333 - 3 * sw - 2 * gap) / 2
    sy0 = 2.0
    for i, (num, title, desc, time, accent, killer) in enumerate(steps):
        x = sx0 + i * (sw + gap)
        y = sy0
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(sw), Inches(sh))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = GOLD if killer else accent
        card.line.width = Pt(3) if killer else Pt(1.5)
        # 大编号
        styled_text(s, x + 0.2, y + 0.1, 1.4, 1.4,
                    num, font=FONT_EN, size=72,
                    color=GOLD if killer else accent, bold=True)
        # 标题
        styled_text(s, x + 0.2, y + 1.5, sw - 0.4, 0.5,
                    title, size=20, color=accent, bold=True)
        # 描述
        styled_text(s, x + 0.2, y + 2.0, sw - 0.4, 0.6,
                    desc, size=13, color=TEXT_SECONDARY)
        # 时间标签
        tb_time = styled_text(s, x + 0.2, y + sh - 0.55, sw - 0.4, 0.4,
                    time, size=14, color=GOLD if killer else accent, bold=True)
        step_cards.append((card, tb_time))

    # 配图:左下 Calcite 联邦图(放大 30% — 视觉锚)
    img_path = PPT_DIR / "images" / "ai" / "p5a-calcite-federation.png"
    pic = _add_picture_with_alpha(s, img_path, 0.2, 5.1, 3.3, 2.4, alpha_pct=100)
    spTree = s.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)

    # 8 动画:1 标题 + 3 步骤卡 + 3 时间标签 + 1 杀手锏脉冲
    add_anim(s, tb_t, "fade_in", delay_ms=0, dur_ms=500)
    for i, (card, tb_time) in enumerate(step_cards):
        add_anim(s, card, "fade_in", delay_ms=400 + i * 300, dur_ms=500)
        add_anim(s, tb_time, "fade_in", delay_ms=400 + i * 300, dur_ms=500)
    # sql-forge-mcp 杀手锏脉冲
    add_anim(s, step_cards[2][0], "pulse", delay_ms=2000, dur_ms=1500, loop=True)


def slide_5b_forge_4plus6(prs):
    """P5-b SQL工坊 · 4 starter + 6 功能 + 3 chase pulse(策略 C)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # 标题
    styled_text(s, 0.5, 0.3, 12.333, 0.5,
                "SQL工坊 · 4 starter + 6 功能",
                size=24, bold=True)

    # 4 starter 横条
    starters = [
        ("sql-forge-starter", SEMANTIC_GREEN),
        ("sql-forge-calcite", JAVA_BLUE),
        ("sql-forge-web", JAVA_BLUE),
        ("sql-forge-mcp", AI_PURPLE),
    ]
    starter_boxes = []
    sw, sh = 3.0, 0.5
    gap = 0.1
    sx0 = (13.333 - 4 * sw - 3 * gap) / 2
    sy = 1.0
    for i, (name, color) in enumerate(starters):
        x = sx0 + i * (sw + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(sy), Inches(sw), Inches(sh))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        styled_text(s, x, sy + 0.05, sw, 0.4,
                    name, font=FONT_MONO, size=11, color=WHITE, bold=True)
        starter_boxes.append(box)

    # 6 功能卡(2x3)— Row1 杀手锏(金边金字),Row2 基础(浅灰边浅蓝字)
    funcs = [
        ("🔌 JSON CRUD ★", GOLD, True),
        ("🌐 Calcite ★", GOLD, True),
        ("🔒 MCP 5 受限 ★", GOLD, True),
        ("🎨 Amis", TEXT_SECONDARY, False),
        ("🏗️ 实体", TEXT_SECONDARY, False),
        ("🛡️ 零直连", TEXT_SECONDARY, False),
    ]
    func_cards = []
    cw, ch = 4.0, 1.4
    gap_x, gap_y = 0.2, 0.2
    cx0 = (13.333 - 3 * cw - 2 * gap_x) / 2
    cy0 = 1.8
    for i, (name, color, killer) in enumerate(funcs):
        col, row = i % 3, i // 3
        x = cx0 + col * (cw + gap_x)
        y = cy0 + row * (ch + gap_y)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(cw), Inches(ch))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        if killer:
            card.line.color.rgb = GOLD
            card.line.width = Pt(3)
        else:
            card.line.color.rgb = DIVIDER
            card.line.width = Pt(1)
        styled_text(s, x, y + 0.4, cw, ch - 0.8,
                    name, size=18, color=color, bold=True)
        func_cards.append(card)

    # 金句
    styled_text(s, 0.5, 6.5, 12.333, 0.5,
                "★ AI 通过 5 个受限通道安全对话业务库,数据不出企业",
                size=18, color=GOLD, bold=True)

    # 11 动画:4 starter + 6 功能 + 3 chase(策略 C)
    for i in range(4):
        add_anim(s, starter_boxes[i], "fade_in",
                 delay_ms=i * 200, dur_ms=500)
    for i in range(6):
        add_anim(s, func_cards[i], "fade_in",
                 delay_ms=1000 + i * 200, dur_ms=500)
    # 3 杀手锏 chase pulse 错开 2s
    add_anim(s, func_cards[0], "pulse", delay_ms=15000, dur_ms=1500, loop=True)  # JSON CRUD
    add_anim(s, func_cards[1], "pulse", delay_ms=17000, dur_ms=1500, loop=True)  # Calcite
    add_anim(s, func_cards[2], "pulse", delay_ms=19000, dur_ms=1500, loop=True)  # MCP


def slide_6_demo1_workflow(prs):
    """P6 演示 1 · 5 步工作流(策略 D:5 步依次入,末步 pulse)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # 标题
    styled_text(s, 0.5, 0.6, 12.333, 0.5,
                 "演示 1 · 一句话出分析报告",
                 size=28, bold=True)

    # 用户原话
    styled_text(s, 0.5, 1.3, 12.333, 0.4,
                 '"上个月各品类销量怎么样?"',
                 size=18, color=TEXT_SECONDARY)

    # 5 步横排
    steps = [
        ("S1", "👤 用户提问", TEXT_SECONDARY),
        ("S2", "🧠 NL2SQL", JAVA_BLUE),
        ("S3", "📊 读表", JAVA_BLUE),
        ("S4", "⚙️ 跨库 JOIN", JAVA_BLUE),
        ("S5", "📄 报告", GOLD),
    ]
    step_boxes = []
    sw, sh = 2.0, 2.5
    gap = 0.3
    total_w = 5 * sw + 4 * gap
    sx0 = (13.333 - total_w) / 2
    sy = 2.5
    for i, (num, name, color) in enumerate(steps):
        x = sx0 + i * (sw + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(sy), Inches(sw), Inches(sh))
        box.fill.solid()
        if color == GOLD:
            box.fill.fore_color.rgb = GOLD_BG
        else:
            box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(3) if color == GOLD else Pt(1)
        # S1-S5 标签
        styled_text(s, x, sy + 0.2, sw, 0.4,
                     num, font=FONT_EN, size=14, color=color, bold=True)
        # 步骤名
        styled_text(s, x, sy + 0.8, sw, 1.5,
                     name, size=16, color=color, bold=True)
        step_boxes.append(box)

    # 箭头(4 个)
    for i in range(4):
        x = sx0 + (i + 1) * sw + i * gap + 0.05
        styled_text(s, x, sy + sh/2 - 0.2, gap - 0.1, 0.4,
                     "→", font=FONT_EN, size=20, color=GOLD, bold=True)

    # 底部金句
    styled_text(s, 0.5, 6.0, 12.333, 0.5,
                 "▼ 接下来看 42 秒真实录屏",
                 size=18, color=AI_PURPLE, bold=True)

    # 6 动画:5 步 0/0.4/0.8/1.2/1.6s 阶梯入场 + 末步 pulse(体现流程而非并列)
    for i in range(5):
        add_anim(s, step_boxes[i], "fade_in",
                 delay_ms=i * 400, dur_ms=500)
    add_anim(s, step_boxes[4], "pulse", delay_ms=2500, dur_ms=1500, loop=True)


def slide_7_demo2_workflow(prs):
    """P7 演示 2 · 5 步工作流"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    styled_text(s, 0.5, 0.6, 12.333, 0.5,
                 "演示 2 · 一句话生成 CRUD",
                 size=28, bold=True)

    styled_text(s, 0.5, 1.3, 12.333, 0.4,
                 '"做一个商品管理 CRUD 页面"',
                 size=18, color=TEXT_SECONDARY)

    steps = [
        ("S1", "👤 用户提问", TEXT_SECONDARY),
        ("S2", "🧠 web.st", AI_PURPLE),
        ("S3", "📊 读表", AI_PURPLE),
        ("S4", "🔧 生成 Amis", AI_PURPLE),
        ("S5", "🌐 URL", GOLD),
    ]
    step_boxes = []
    sw, sh = 2.0, 2.5
    gap = 0.3
    total_w = 5 * sw + 4 * gap
    sx0 = (13.333 - total_w) / 2
    sy = 2.5
    for i, (num, name, color) in enumerate(steps):
        x = sx0 + i * (sw + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(sy), Inches(sw), Inches(sh))
        box.fill.solid()
        if color == GOLD:
            box.fill.fore_color.rgb = GOLD_BG
        else:
            box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(3) if color == GOLD else Pt(1)
        styled_text(s, x, sy + 0.2, sw, 0.4,
                     num, font=FONT_EN, size=14, color=color, bold=True)
        styled_text(s, x, sy + 0.8, sw, 1.5,
                     name, size=16, color=color, bold=True)
        step_boxes.append(box)

    for i in range(4):
        x = sx0 + (i + 1) * sw + i * gap + 0.05
        styled_text(s, x, sy + sh/2 - 0.2, gap - 0.1, 0.4,
                     "→", font=FONT_EN, size=20, color=GOLD, bold=True)

    styled_text(s, 0.5, 6.0, 12.333, 0.5,
                 "▼ 接下来看 40 秒真实录屏",
                 size=18, color=AI_PURPLE, bold=True)

    for i in range(5):
        add_anim(s, step_boxes[i], "fade_in",
                 delay_ms=i * 400, dur_ms=500)
    add_anim(s, step_boxes[4], "pulse", delay_ms=2500, dur_ms=1500, loop=True)


def slide_8_compare(prs):
    """P8 实战对比 · 左技术派 ✓ + 中 3 柱状图(传统/集成AI/JavaBrain) + 右商业派 5 项 + 杀手锏金脉冲(策略 C:13 动画,精简分组滑入 3.5s)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # 标题
    styled_text(s, 0.5, 0.3, 12.333, 0.5,
                 "实战测试 · 7 张表 / 0 漂移 / 节省 >80%",
                 size=24, bold=True)

    # 左半:技术派 4 行 ✓
    styled_text(s, 0.3, 0.95, 4.0, 0.35,
                 "技术派 · 4 项测试",
                 size=15, color=JAVA_BLUE, bold=True)

    items_tech = [
        "✓ 7 张表测试(简单/字典/自关联)",
        "✓ 0 漂移 / 0 错误链接 / 0 乱码",
        "✓ 5/5 AI 字段推断自检通过",
        "✓ 节省 CRUD 开发时间 >80%",
    ]
    tech_boxes = []
    for i, txt in enumerate(items_tech):
        is_killer = ">80%" in txt
        tb = styled_text(s, 0.4, 1.5 + i * 0.55, 4.0, 0.4,
                          txt, size=13,
                          color=GOLD if is_killer else TEXT_PRIMARY,
                          bold=is_killer)
        tech_boxes.append(tb)

    # 中部:3 柱状图(传统 vs 集成 AI vs JavaBrain)
    styled_text(s, 4.7, 0.95, 3.5, 0.35,
                 "开发时间 · 3 种方案",
                 size=15, color=AI_PURPLE, bold=True)

    # 柱状图数据
    bar_data = [
        ("传统开发", 30, DIVIDER),
        ("集成 AI", 10, SEMANTIC_RED),
        ("JavaBrain", 1, SEMANTIC_GREEN),
    ]
    bar_shapes = []
    bar_w = 0.9
    bar_gap = 0.25
    bar_x0 = 4.9
    bar_baseline_y = 5.0  # 底部基线
    bar_max_h = 2.8  # 30天 对应 2.8 in
    bar_label_y = 5.15
    bar_title_y = 5.55
    for i, (label, days, color) in enumerate(bar_data):
        h = (days / 30.0) * bar_max_h
        x = bar_x0 + i * (bar_w + bar_gap)
        y = bar_baseline_y - h
        # 柱子
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y),
                                  Inches(bar_w), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        # 数字标签(柱子顶部)
        styled_text(s, x, y - 0.35, bar_w, 0.3,
                    f"{days}天", size=11, color=color, bold=True, center=True)
        # 底部标签
        styled_text(s, x, bar_label_y, bar_w, 0.3,
                    label, size=11, color=TEXT_PRIMARY, bold=True, center=True)
        bar_shapes.append(bar)

    # 右半:商业派 5 行对比表
    styled_text(s, 8.3, 0.95, 4.8, 0.35,
                 "商业派 · 5 项对比",
                 size=15, color=AI_PURPLE, bold=True)

    rows_biz = ["业务取数", "CRUD 页面", "跨库 JOIN", "私有化", "列名识别"]
    biz_boxes = []
    for i, txt in enumerate(rows_biz):
        # 行背景
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(8.3), Inches(1.5 + i * 0.5),
                                  Inches(4.8), Inches(0.45))
        bg.fill.solid()
        bg.fill.fore_color.rgb = GREEN_BG if i % 2 == 0 else WHITE
        bg.line.color.rgb = DIVIDER
        bg.line.width = Pt(1)
        # 行文字
        row_tb = styled_text(s, 8.5, 1.55 + i * 0.5, 3.0, 0.4,
                              txt, size=14, center=False)
        # ✓
        styled_text(s, 11.7, 1.55 + i * 0.5, 1.2, 0.4,
                     "✓", font=FONT_EN, size=16, color=SEMANTIC_GREEN, bold=True)
        biz_boxes.append(row_tb)

    # 底部金句
    styled_text(s, 0.5, 6.3, 12.333, 0.5,
                 "★ JavaBrain 在 安全 + 智能 + 私有化 三角都做到",
                 size=18, color=GOLD, bold=True)

    # 13 动画:策略 C(精简分组滑入 — 整组 3.5s 内完成,3 段紧贴)
    for i in range(4):                                              # 技术派 4 项 0-900ms
        add_anim(s, tech_boxes[i], "fade_in",
                 delay_ms=i * 225, dur_ms=400)
    for i, bar in enumerate(bar_shapes):                            # 柱状图 3 项 1200-1800ms
        add_anim(s, bar, "fade_in", delay_ms=1200 + i * 200, dur_ms=400)
    for i in range(5):                                              # 商业派 5 项 2100-2900ms
        add_anim(s, biz_boxes[i], "fade_in",
                 delay_ms=2100 + i * 200, dur_ms=400)
    add_anim(s, tech_boxes[3], "pulse", delay_ms=12000, dur_ms=1500, loop=True)


def slide_9_roadmap(prs):
    """P9 路线图 · 双轨时间线 + 2 列未来目标(灵梭 60% 主 / SQL工坊 40% 辅)+ 双引擎 3 金钩(13 动画)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # 标题
    tb_t = styled_text(s, 0.5, 0.3, 12.333, 0.6,
                       "JavaBrain 路线图",
                       size=28, bold=True)

    # 双轨时间线背景图(全宽,中部 100% alpha)
    img_path = PPT_DIR / "images" / "ai" / "p9_dual_timeline.png"
    pic = _add_picture_with_alpha(s, img_path, 0.5, 0.95, 12.333, 1.2, alpha_pct=100)
    spTree = s.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)

    # 左列:灵梭 (Java 蓝,60% 宽 = 7.5)
    loom_items = [
        ("1. Agentic Workflow", "智能体工作流:自动拆解 + 多步推理 + 反思重试"),
        ("2. Graph-RAG + 实时认知", "知识图谱融合 + 毫秒级增量更新"),
        ("3. 多模态交互", "语音 / 图像 / AI 悬浮球"),
        ("4. Skill 插件生态市场", "标准化打包 + 私有 Skill 市场"),
    ]
    styled_text(s, 0.3, 2.4, 7.6, 0.4,
                "灵梭 · AI 赋能中心 (主线 60%)",
                size=18, color=JAVA_BLUE, bold=True)
    loom_boxes = []
    for i, (title, desc) in enumerate(loom_items):
        y = 2.95 + i * 0.6
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.3), Inches(y),
                                    Inches(7.6), Inches(0.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = JAVA_BLUE
        card.line.width = Pt(1.5)
        styled_text(s, 0.5, y + 0.04, 7.4, 0.25,
                    title, size=13, color=JAVA_BLUE, bold=True)
        styled_text(s, 0.5, y + 0.28, 7.4, 0.22,
                    desc, size=11, color=TEXT_SECONDARY)
        loom_boxes.append(card)

    # 右列:SQL工坊 (AI 紫,40% 宽 = 5.0)
    forge_items = [
        ("1. Text-to-SQL 终极形态", "准确率 >95% · 本地代码大模型"),
        ("2. HTAP 湖仓一体", "Flink + Iceberg/Hudi"),
        ("3. 智能数据治理", "血缘 + 脱敏 + 审计"),
        ("4. Serverless 弹性", "计算/存储分离"),
    ]
    styled_text(s, 8.1, 2.4, 5.0, 0.4,
                "SQL工坊 · 数据引擎",
                size=16, color=AI_PURPLE, bold=True)
    forge_boxes = []
    for i, (title, desc) in enumerate(forge_items):
        y = 2.95 + i * 0.6
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(8.1), Inches(y),
                                    Inches(5.0), Inches(0.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = AI_PURPLE
        card.line.width = Pt(1.5)
        styled_text(s, 8.3, y + 0.04, 4.8, 0.25,
                    title, size=12, color=AI_PURPLE, bold=True)
        styled_text(s, 8.3, y + 0.28, 4.8, 0.22,
                    desc, size=10, color=TEXT_SECONDARY)
        forge_boxes.append(card)

    # 底部 3 金钩:双引擎协同愿景
    hooks = [
        "★ 让 AI 懂数据,让数据懂 AI",
        "★ 终极形态:对话式数据分析 (CDA) 标杆",
        "★ 零代码 AI 应用工厂 (Amis + Skill 绑定)",
    ]
    hook_boxes = []
    for i, txt in enumerate(hooks):
        y = 5.6 + i * 0.42
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(2.0), Inches(y),
                                  Inches(9.333), Inches(0.35))
        box.fill.solid()
        box.fill.fore_color.rgb = GOLD_BG
        box.line.color.rgb = GOLD
        box.line.width = Pt(1.5)
        styled_text(s, 2.0, y + 0.02, 9.333, 0.33,
                    txt, size=15, color=GOLD, bold=True, center=True)
        hook_boxes.append(box)

    # 13 动画:1 标题 + 4 灵梭 + 4 SQL工坊 + 3 金钩 + 1 金钩[1] 脉冲
    add_anim(s, tb_t, "fade_in", delay_ms=0, dur_ms=500)
    for i, box in enumerate(loom_boxes):
        add_anim(s, box, "fade_in", delay_ms=400 + i * 200, dur_ms=500)
    for i, box in enumerate(forge_boxes):
        add_anim(s, box, "fade_in", delay_ms=1200 + i * 200, dur_ms=500)
    for i, box in enumerate(hook_boxes):
        add_anim(s, box, "fade_in", delay_ms=2000 + i * 300, dur_ms=500)
    add_anim(s, hook_boxes[1], "pulse", delay_ms=3500, dur_ms=1500, loop=True)


def slide_10_ending(prs):
    """P10 结尾 · 3 句话 + 金句 + 启动日志 + 仓库地址(策略 B:4 入场 + 2 循环)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_LIGHT

    # JavaBrain 大字(标题上方大脑小图标作为 logo,透明背景版避免被裁)+ 8s 慢转
    img_path = PPT_DIR / "images" / "ai" / "p10_brain_logo.png"
    p10_logo = s.shapes.add_picture(str(img_path), Inches(5.9), Inches(0.05),
                                     width=Inches(1.5), height=Inches(0.6))

    # JavaBrain 大字(下移以避开 logo)
    tb_t = styled_text(s, 0.5, 0.75, 12.333, 1.2,
                        "JavaBrain",
                        font=FONT_EN, size=60, color=TEXT_PRIMARY, bold=True)

    # 3 句话
    sentences = [
        "1. 开箱即用 + 私有化 + 业务语义",
        "2. 一行命令启动,一个文件改 AI 行为",
        "3. 让每个 Spring Boot 项目都自带 AI 能力",
    ]
    sent_boxes = []
    for i, txt in enumerate(sentences):
        tb = styled_text(s, 0.5, 2.0 + i * 0.5, 12.333, 0.4,
                          txt, size=18)
        sent_boxes.append(tb)

    # 金句
    killer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(2.5), Inches(4.2),
                                  Inches(8.333), Inches(0.7))
    killer.fill.solid()
    killer.fill.fore_color.rgb = GOLD_BG
    killer.line.color.rgb = GOLD
    killer.line.width = Pt(3)
    styled_text(s, 2.5, 4.25, 8.333, 0.6,
                 "★ 让 AI 懂数据,让数据懂 AI — 打造对话式数据分析标杆",
                 size=20, color=GOLD, bold=True)

    # 终端框
    term = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(2.5), Inches(5.2),
                                Inches(8.333), Inches(0.6))
    term.fill.solid()
    term.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    term.line.fill.background()
    styled_text(s, 2.5, 5.25, 8.333, 0.5,
                 "$ Started LoomAgentApplication in 25.394 seconds",
                 font=FONT_MONO, size=14, color=SEMANTIC_GREEN)

    # 仓库地址
    styled_text(s, 0.5, 6.1, 12.333, 0.5,
                 "github.com/wb04307201/spring-ai-loom-agent  ·  gitee.com/wb04307201/spring-ai-loom-agent",
                 font=FONT_MONO, size=12, color=TEXT_SECONDARY)

    # 致谢(开源项目 — 用通用联系信息,留出占位供演讲者替换)
    styled_text(s, 0.5, 6.7, 12.333, 0.3,
                 "感谢聆听 · Star 支持: github.com/wb04307201/spring-ai-loom-agent",
                 size=12, color=TEXT_SECONDARY)

    # 7 动画:logo 旋转 + 5 入场 + 1 金句循环
    add_anim(s, p10_logo, "spin", delay_ms=0, dur_ms=8000, loop=True)  # 8s 慢转
    add_anim(s, tb_t, "fade_in", delay_ms=0, dur_ms=500)
    add_anim(s, sent_boxes[0], "fade_in", delay_ms=600, dur_ms=500)
    add_anim(s, sent_boxes[1], "fade_in", delay_ms=1200, dur_ms=500)
    add_anim(s, sent_boxes[2], "fade_in", delay_ms=1800, dur_ms=500)
    add_anim(s, killer, "pulse", delay_ms=6000, dur_ms=1500, loop=True)  # 金句持续脉冲
    add_anim(s, term, "fade_in", delay_ms=4000, dur_ms=500)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_1_cover(prs)
    slide_2_pain(prs)
    slide_3_position(prs)
    slide_4a_loom_intro(prs)
    slide_4b_loom_modules(prs)
    slide_5a_forge_intro(prs)
    slide_5b_forge_4plus6(prs)
    slide_6_demo1_workflow(prs)
    slide_7_demo2_workflow(prs)
    slide_8_compare(prs)
    slide_9_roadmap(prs)
    slide_10_ending(prs)
    # 统一加页码(右下角,"XX / 12"格式)
    for i, slide in enumerate(prs.slides, 1):
        add_page_number(slide, i, total=12)
    prs.save(str(OUTPUT))
    print(f"OK: {OUTPUT} (12 pages)")


if __name__ == "__main__":
    main()
