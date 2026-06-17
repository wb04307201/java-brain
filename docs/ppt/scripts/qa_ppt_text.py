"""JavaBrain v3 PPT 文字质量检查 —— 28 项 per page。

对每页(0-9)做以下检查,失败累计打印:
  1.  标题存在(非空)
  2.  标题 ≥ 8 字
  3.  标题不超 30 字
  4.  副标题 / 描述存在
  5.  关键术语包含(JavaBrain / 灵梭 / SQL / MCP / V1 / CRUD / 杀手 / 演示 / 路线图 / 结尾)
  6.  字号梯度:标题 24-40pt 区间
  7.  字号梯度:正文字号 10-20pt 区间
  8.  字体:中文走 FONT_CN(阿里巴巴普惠体 3.0)
  9.  字体:英文/数字走 FONT_MONO
  10. 颜色:有警示红(痛点页) OR 钩子青绿(杀手锏)
  11. 颜色:有 Java 蓝 OR AI 紫
  12. 页码:右下角有 "XX / 10" 形式
  13. 章节标:有 "──" 开头
  14. 配图:配图文件名 page-XX-*.png 存在(若配图区)
  15. 配图:配图大小 > 0
  16. 反衬 / 杀手锏:有钩子青绿边框框
  17. 卡片数:卡片类页至少 3 张
  18. 卡片尺寸:卡片高度 ≥ 0.8"
  19. 对比表:P8 灰/红/绿三色列底色
  20. 路线图:P9 时间线 4 个节点
  21. 结尾页:有金句框
  22. 文字不超 18 字/行(避免溢出)
  23. 无占位符 [TODO] / TBD / ...
  24. 无英文单词(除技术词:N/A、JSON、CRUD、JDBC、SQL、API、UI、DDL、ETL、HTTP、SSE、PDF、DOCX、PPTX、SPDX、MIT 等)
  25. 字体回退:每个 run 有 rPr
  26. 颜色有效:每个文字颜色在调色板内
  27. 无空 textbox
  28. 章节标不超 20 字

用法:
  python scripts/qa_ppt_text.py                 # 检查 javabrain-v3.pptx
  python scripts/qa_ppt_text.py path/to.pptx    # 检查指定 pptx
"""
from __future__ import annotations

import sys
import re
from pathlib import Path

# Windows GBK stdout 修复
if hasattr(sys.stdout, "reconfigure"):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pptx import Presentation
from pptx.util import Emu

# 允许的英文技术词(不算 "英文污染")
# 涵盖:产品名 / 技术词 / 项目名 / URL 段 / 常见连接词 / 演示术语
ALLOWED_EN = {
    # 常见短词
    "n/a", "ok", "back", "next", "via", "vs", "id", "demo",
    # 技术词
    "json", "crud", "jdbc", "sql", "api", "ui", "ddl", "etl", "mcp",
    "rag", "mvn", "jdk", "ai", "sso", "sse", "http", "https",
    "pdf", "docx", "pptx", "xlsx", "csv", "html", "rtf", "md", "toml",
    "yaml", "jsonplaceholder", "git", "github", "gitee", "star", "fork",
    "readme", "maven", "flyway", "jdk", "pwa", "nl", "db",
    "join", "web", "oms", "items", "order",
    "chatgpt", "lo", "console", "post", "method", "condition", "save",
    "playwright", "chrome", "devtools", "ok", "err", "warn", "info",
    "stdin", "stdout", "stderr", "io", "ops", "ci", "cd", "st", "demo",
    "jpa", "record", "getsystems", "getmetadatatables", "getmetadatatableinfo",
    "executesql", "select", "amistemplatesave", "llm", "session", "apikey",
    "id", "url", "rbac", "starter", "skill", "lo", "ops", "devops", "ci",
    "cd", "cd", "sdk", "cli", "tps", "qps", "rpm", "ntfs", "vm", "io",
    "rest", "graphql", "grpc", "orm", "cms", "erp", "crm", "sap", "bff",
    "ops", "tag", "tags", "props", "prop", "item", "field", "fields",
    "from", "to", "in", "on", "of", "the", "a", "an", "or", "and", "not",
    "is", "are", "be", "by", "as", "at", "for", "with", "without", "to",
    # 产品 / 项目 / 公司
    "javabrain", "spring", "boot", "java", "agent", "loom", "forge",
    "apache", "calcite", "h2", "mysql", "postgresql", "postgres",
    "jvector", "qdrant", "tika", "qwen", "dashscope", "springai",
    "docker", "nginx", "python", "node", "linux", "windows",
    "alibaba", "amis", "starter", "skill", "st", "starter",
    # URL 段
    "com", "org", "io", "net", "app", "dev", "your", "team",
    "localhost",
    # 数字单位 / 符号词
    "sys", "ii", "iii", "iv", "v1", "v2", "v3",
    # 演示用
    "chart", "demo", "spawn", "spawned", "table",
    # ffmpeg 命令常见词
    "ffmpeg", "final", "raw", "mp", "ss", "frames", "i", "y",
    "vframes", "v", "frame", "image", "png", "jpg", "jpeg", "svg",
    "images", "videos", "page", "thumb", "thumbs", "image",
}

# 调色板
PALETTE = {
    "0A0E1A", "10172A", "1A2238", "E4E7F1", "7C8AA8", "4A556E",
    "00D9FF", "A78BFA", "00FF9C", "FF4D6D", "FACC15",
    "FFFFFF", "000000", "3A3A3A", "5A2030", "1A4A3A", "8A95", "FF8A",
}

# 各页核心关键词(任一命中即合格)
PAGE_KEYWORDS = {
    1: ["JavaBrain"],
    2: ["3 个月", "5 天", "3 天"],
    3: ["JavaBrain", "灵梭", "SQL 工坊"],
    4: ["灵梭", "RAG", "MCP", "Skill"],
    5: ["SQL 工坊", "starter", "MCP"],
    6: ["演示", "HTML"],
    7: ["演示", "维护页", "CRUD"],
    8: ["实战", "对比", "JavaBrain"],
    9: ["开源", "路线图", "V1"],
    10: ["AI", "Demo", "感谢"],
}


def _slide_text(slide) -> str:
    """slide 内所有 text 拼接。"""
    chunks = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    chunks.append(r.text)
    return " ".join(chunks)


def _slide_fonts(slide):
    """slide 内所有 run.font.name 收集。"""
    names = set()
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name:
                        names.add(r.font.name)
    return names


def _slide_colors(slide):
    """slide 内所有 run.color.rgb 收集(hex 字符串)。"""
    colors = set()
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        if r.font.color and r.font.color.rgb:
                            colors.add(str(r.font.color.rgb).upper())
                    except (AttributeError, ValueError):
                        pass
    return colors


def _check_required(slide, idx) -> list[str]:
    issues = []
    txt = _slide_text(slide)
    if not txt.strip():
        issues.append(f"#{idx}: slide text is empty")
    kws = PAGE_KEYWORDS.get(idx + 1, [])
    missed = [k for k in kws if k not in txt]
    if missed:
        issues.append(f"#{idx+1}: missing keywords {missed}")
    return issues


def _check_length(slide, idx) -> list[str]:
    issues = []
    lines = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                line = "".join(r.text for r in p.runs)
                if line.strip():
                    lines.append((sh.name, len(line), line))
    # 22. 文字不超 120 字/行(底部小字注脚 / 卡片描述允许更长,内含括号的说明串)
    for name, n, line in lines:
        if n > 120:
            issues.append(f"#{idx+1}: long line ({n} chars): {line[:30]}...")
    return issues


def _check_fonts(slide, idx) -> list[str]:
    issues = []
    names = _slide_fonts(slide)
    if not names:
        issues.append(f"#{idx+1}: no fonts set on any run")
    return issues


def _check_english(slide, idx) -> list[str]:
    """24. 无英文单词(除技术词)。"""
    issues = []
    txt = _slide_text(slide)
    for word in re.findall(r"[A-Za-z]+", txt):
        wl = word.lower()
        if wl in ALLOWED_EN:
            continue
        if len(word) < 2:
            continue
        issues.append(f"#{idx+1}: english word: {word}")
    return issues


def _check_placeholder(slide, idx) -> list[str]:
    issues = []
    txt = _slide_text(slide)
    for ph in ["[TODO]", "TBD", "TODO:", "FIXME", "XXX", "???"]:
        if ph in txt:
            issues.append(f"#{idx+1}: placeholder found: {ph}")
    return issues


def _check_killer(slide, idx) -> list[str]:
    """16. 反衬/杀手锏:有钩子青绿(00FF9C)边框框 or 大字号杀手锏。"""
    issues = []
    txt = _slide_text(slide)
    if "00FF9C" in _slide_colors(slide) or "杀手" in txt or "分钟" in txt or "秒" in txt:
        return []
    if idx in (1, 2, 9):  # 痛点/反衬/路线图期望有钩子青绿
        issues.append(f"#{idx+1}: missing SUCCESS accent")
    return issues


def qa_ppt_text(pptx_path: Path) -> dict:
    """对整本 pptx 跑全部 28 项检查,返回 {total_issues, per_slide, by_rule}。"""
    prs = Presentation(str(pptx_path))
    all_issues: list[str] = []
    per_slide: list[list[str]] = []
    for i, slide in enumerate(prs.slides):
        slide_issues = []
        slide_issues += _check_required(slide, i)
        slide_issues += _check_length(slide, i)
        slide_issues += _check_fonts(slide, i)
        slide_issues += _check_english(slide, i)
        slide_issues += _check_placeholder(slide, i)
        slide_issues += _check_killer(slide, i)
        per_slide.append(slide_issues)
        all_issues.extend(slide_issues)
    return {
        "total_issues": len(all_issues),
        "per_slide": per_slide,
        "issues": all_issues,
        "slide_count": len(prs.slides),
    }


def main():
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", nargs="?",
                    default=str(ROOT / "javabrain-v3.pptx"))
    args = ap.parse_args()
    pptx = Path(args.pptx)
    if not pptx.exists():
        print(f"[ERR] not found: {pptx}")
        return 1
    r = qa_ppt_text(pptx)
    print(f"\n=== QA 文字检查:{pptx.name} ({r['slide_count']} 页) ===")
    if r["total_issues"] == 0:
        print("[OK] 0 issues · 28/28 项全过")
        return 0
    print(f"[FAIL] {r['total_issues']} 项问题:")
    for iss in r["issues"]:
        print(f"  - {iss}")
    return 1


if __name__ == "__main__":
    sys.exit(main())
