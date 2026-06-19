#!/usr/bin/env python3
"""为 P6/P7 添加 AI 工作步骤展示内容。

参考 javabrain_v1.pptx 的结构，添加:
- 标题
- 演示背景标注
- 演示语
- 画面步骤分解
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 从 tools.py 导入颜色常量
sys.path.insert(0, str(Path(__file__).parent))
from tools import (
    BG_DEEP,
    BG_PANEL,
    TEXT_PRIMARY,
    TEXT_SECONDARY,
    PRIMARY,
    SUCCESS,
    FONT_CN,
    FONT_MONO,
)


def add_text_box(slide, x, y, w, h, text, *, font=FONT_MONO, size=14,
                 color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT):
    """添加文本框。"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return txBox


def add_demo_page_6(slide):
    """添加 P6 录屏 1 的内容。"""
    # 标题
    add_text_box(slide, 0.75, 0.22, 10.5, 0.5,
                 "录屏 1:一句话出分析报告",
                 font=FONT_CN, size=20, color=TEXT_PRIMARY, bold=True)

    # 演示背景标注
    add_text_box(slide, 0.75, 0.70, 10.5, 0.35,
                 "演示基于 oms 订单管理系统(用 SQL 工坊搭建)",
                 font=FONT_CN, size=11, color=TEXT_SECONDARY)

    # 演示语
    add_text_box(slide, 0.75, 1.18, 12.0, 0.5,
                 "演示语:订单系统各分类商品数,画柱状图,保存 HTML 报告",
                 font=FONT_CN, size=14, color=PRIMARY)

    # 画面步骤分解标题
    add_text_box(slide, 0.75, 1.85, 12.0, 0.4,
                 "画面 6 步分解:",
                 font=FONT_CN, size=12, color=TEXT_SECONDARY, bold=True)

    # 步骤
    steps = [
        ("① 输入", "5s"),
        ("② AI 推理", "15s"),
        ("③ 表格 + 柱状图", "20s"),
        ("④ 保存 HTML", "10s"),
        ("⑤ 打开预览", "15s"),
        ("⑥ 返回聊天", "5s"),
    ]

    for idx, (step, time) in enumerate(steps):
        y = 2.35 + idx * 0.55
        # 步骤名
        add_text_box(slide, 1.0, y, 6.0, 0.4,
                     step,
                     font=FONT_CN, size=13, color=TEXT_PRIMARY)
        # 时间
        add_text_box(slide, 7.5, y, 2.0, 0.4,
                     time,
                     font=FONT_MONO, size=13, color=SUCCESS, align=PP_ALIGN.RIGHT)


def add_demo_page_7(slide):
    """添加 P7 录屏 2 的内容。"""
    # 标题
    add_text_box(slide, 0.75, 0.22, 10.5, 0.5,
                 "录屏 2:一句话生成 CRUD",
                 font=FONT_CN, size=20, color=TEXT_PRIMARY, bold=True)

    # 演示背景标注
    add_text_box(slide, 0.75, 0.70, 10.5, 0.35,
                 "演示基于 oms 订单管理系统(用 SQL 工坊搭建)",
                 font=FONT_CN, size=11, color=TEXT_SECONDARY)

    # 演示语
    add_text_box(slide, 0.75, 1.18, 12.0, 0.5,
                 "演示语 2 轮：① 帮我做一个订单明细 (ORDER_ITEMS) 的维护页  ② 没问题，直接生成",
                 font=FONT_CN, size=14, color=PRIMARY)

    # 画面步骤分解标题
    add_text_box(slide, 0.75, 1.85, 12.0, 0.4,
                 "画面 6 步分解:",
                 font=FONT_CN, size=12, color=TEXT_SECONDARY, bold=True)

    # 步骤
    steps = [
        ("① 输入 + 推断", "15s"),
        ("② 确认生成", "10s"),
        ("③ 3 行输出", "5s"),
        ("④ 切 oms 登录", "20s"),
        ("⑤ 演示功能", "30s"),
        ("⑥ 切回聊天", "5s"),
    ]

    for idx, (step, time) in enumerate(steps):
        y = 2.35 + idx * 0.55
        # 步骤名
        add_text_box(slide, 1.0, y, 6.0, 0.4,
                     step,
                     font=FONT_CN, size=13, color=TEXT_PRIMARY)
        # 时间
        add_text_box(slide, 7.5, y, 2.0, 0.4,
                     time,
                     font=FONT_MONO, size=13, color=SUCCESS, align=PP_ALIGN.RIGHT)


def main():
    """主流程。"""
    script_dir = Path(__file__).parent
    ppt_dir = script_dir.parent
    ppt_path = ppt_dir / "javabrain_v2.pptx"

    print(f"打开 PPT: {ppt_path}")
    prs = Presentation(str(ppt_path))

    # P6 (索引 5)
    if len(prs.slides) >= 6:
        slide6 = prs.slides[5]
        print("添加 P6 内容...")
        add_demo_page_6(slide6)

    # P7 (索引 6)
    if len(prs.slides) >= 7:
        slide7 = prs.slides[6]
        print("添加 P7 内容...")
        add_demo_page_7(slide7)

    # 保存
    prs.save(str(ppt_path))
    print(f"\n[OK] PPT 已更新：{ppt_path}")

    return 0


if __name__ == "__main__":
    sys.exit(main())
