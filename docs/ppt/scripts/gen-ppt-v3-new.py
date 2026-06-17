"""JavaBrain v3 PPT 生成器 v2(10 页)—— 严格按 v1 版式:深色背景 + 大图占主 + 紧凑文字

版式骨架(继承自 javabrain_v1.pptx):
  - 每页固定结构:左上 0.12x0.70 章节色条 + 28pt 冷白标题 + 12pt 冷灰副
  - P1 全屏背景图 + 底部黑条标题
  - P2-P5 / P9 顶部标题副 → 中央大图占主 → 底部深底反衬条(18pt 橙字 + 10pt 灰注脚)
  - P6/P7 演示页:演示语框 + 6 步流程色块(无时间戳)
  - P8 左实战数据(>80% 橙大字) + 右 4×5 对比表(灰/红/绿/深撞色块)
  - P10 顶部 3 句(深底) + 中央配图 + 底部黑条橙金句

内容来自 OUTLINE_v3.md。

用法:python scripts/gen-ppt-v3-new.py
"""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import tools as T
from tools import (
    FONT_CN, FONT_MONO,
    BG_DEEP, BG_PANEL, GRID_LINE, DIVIDER,
    TEXT_PRIMARY, TEXT_SECONDARY, TEXT_DIM,
    PRIMARY, ACCENT, SUCCESS, DANGER, WARN,
    add_text, add_rect, add_picture_with_alpha, new_presentation,
)

IMG = ROOT / "images"
OUT = ROOT / "javabrain-v3.pptx"
TOTAL = 10

# v1 配色(深色 HUD 同色系,文字色偏冷白)
SLATE_BG = T.RGBColor(0x0F, 0x17, 0x2A)        # 深蓝黑主底
DEEP_PANEL = T.RGBColor(0x1E, 0x2A, 0x47)       # 深色块(底部反衬)
DEEPER = T.RGBColor(0x07, 0x0C, 0x19)           # 更深色块(P6/P7 占位区 / P10 黑条)
DARKER = T.RGBColor(0x1A, 0x1F, 0x2E)           # 对比表数据行
COOL_WHITE = T.RGBColor(0xF8, 0xFA, 0xFC)       # 冷白主字
COOL_GRAY = T.RGBColor(0xCB, 0xD5, 0xE1)        # 冷灰副字
COOL_DIM = T.RGBColor(0x94, 0xA3, 0xB8)         # 注脚灰
ORANGE = T.RGBColor(0xF5, 0x9E, 0x0B)           # 反衬橙(v1 招牌色)
GREEN = T.RGBColor(0x10, 0xB9, 0x81)            # 钩子绿
RED = T.RGBColor(0xEF, 0x44, 0x44)              # 警示红
PURPLE = T.RGBColor(0x6B, 0x46, 0xC1)           # 紫
JAVA_BLUE_DARK = T.RGBColor(0x00, 0x73, 0x96)   # 蓝(对比表头)


# ====================================================================
# v1 版式公共组件
# ====================================================================
def v1_title(slide, *, accent: T.RGBColor, title: str, sub: str) -> None:
    """左上章节色条 + 28pt 冷白标题 + 12pt 冷灰副。"""
    add_rect(slide, 0.50, 0.32, 0.12, 0.70,
             fill=accent, line_color=None, line_width_pt=0)
    add_text(slide, 0.75, 0.22, 10.50, 0.50, title,
             font=FONT_CN, size=28, color=COOL_WHITE, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(slide, 0.75, 0.70, 10.50, 0.35, sub,
             font=FONT_CN, size=12, color=COOL_GRAY,
             align=PP_ALIGN.LEFT)


def v1_bottom_bar(slide, *, killer: str, note: str) -> None:
    """底部反衬条:深底 + 18pt 橙杀手锏 + 10pt 灰注脚。"""
    add_rect(slide, 0.60, 6.30, 12.10, 0.55,
             fill=DEEP_PANEL, line_color=None, line_width_pt=0)
    add_text(slide, 0.70, 6.32, 11.90, 0.50, killer,
             font=FONT_CN, size=18, color=ORANGE, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(slide, 0.60, 6.95, 12.10, 0.40, note,
             font=FONT_CN, size=10, color=COOL_DIM,
             align=PP_ALIGN.LEFT)


def v1_page_bg(slide, color=SLATE_BG) -> None:
    """slide 背景色。"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


# ====================================================================
# P1 封面(全屏背景图 + 底部黑条标题)
# ====================================================================
def slide_1_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    img = IMG / "page-01-cover.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(0), Inches(0),
                             Inches(13.33), Inches(7.5))
    # 底部黑条
    add_rect(s, 0.00, 6.00, 13.33, 1.50,
             fill=SLATE_BG, line_color=None, line_width_pt=0)
    add_text(s, 0.80, 6.10, 11.70, 0.70, "JavaBrain",
             font=FONT_CN, size=46, color=COOL_WHITE, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, 0.80, 6.75, 11.70, 0.40,
             "Spring Boot AI Agent + 数据库低代码一站式解决方案",
             font=FONT_CN, size=16, color=COOL_GRAY,
             align=PP_ALIGN.LEFT)
    add_text(s, 0.80, 7.10, 11.70, 0.35,
             "灵梭 + SQL 工坊 + SQL 工坊 MCP   ·   开源 · 私有化 · 业务语义",
             font=FONT_CN, size=11, color=COOL_DIM,
             align=PP_ALIGN.LEFT)


# ====================================================================
# P2 痛点
# ====================================================================
def slide_2_pain(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    v1_title(s, accent=DANGER,
             title="企业 AI 落地的 3 个真实数字",
             sub="3 个月 / 5 天 / 3 天  →  3 分钟 / 90 秒 / 10 分钟")
    img = IMG / "page-02-pain.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(2.29), Inches(1.20),
                             Inches(8.71), Inches(4.90))
    # 在底部反衬条杀手锏前加一行警示红数字反衬
    add_text(s, 0.70, 5.80, 11.90, 0.40,
             "★ 3 个月 / 5 天 / 3 天   ←  JavaBrain 压到 3 分钟 / 90 秒 / 10 分钟",
             font=FONT_CN, size=14, color=DANGER, bold=True,
             align=PP_ALIGN.LEFT)
    v1_bottom_bar(s,
                  killer="3 分钟启动  ·  90 秒出报告  ·  10 分钟出可用的页面",
                  note="JavaBrain 把这 3 个数字,从月压到分")


# ====================================================================
# P3 定位 + 组成
# ====================================================================
def slide_3_position(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    v1_title(s, accent=PURPLE,
             title="定位 + 组成",
             sub="三个 Spring Boot 依赖库,按需引入,不需要全栈改造")
    # 核心金句(中间一行)
    add_text(s, 0.60, 1.15, 12.10, 0.50,
             "JavaBrain = 灵梭 AI Agent  +  SQL 工坊  +  SQL 工坊 MCP",
             font=FONT_CN, size=22, color=SUCCESS, bold=True,
             align=PP_ALIGN.CENTER)
    img = IMG / "page-03-position.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(2.83), Inches(1.75),
                             Inches(7.64), Inches(4.30))
    v1_bottom_bar(s,
                  killer="三个都是依赖库,组合 = 企业 AI 落地完整闭环",
                  note="灵梭 / SQL 工坊 / SQL 工坊 MCP,各自独立可用,按需组合")


# ====================================================================
# P4 灵梭 6 大功能
# ====================================================================
def slide_4_loom(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    v1_title(s, accent=PURPLE,
             title="灵梭 · 6 大功能模块",
             sub="Spring AI 编排 + RAG / MCP / Skill 全部装进一个依赖库")
    # 右上角徽章
    add_rect(s, 11.40, 0.40, 1.70, 0.40,
             fill=PURPLE, line_color=None, line_width_pt=0)
    add_text(s, 11.40, 0.40, 1.70, 0.40, "🧩 独立依赖库",
             font=FONT_CN, size=11, color=COOL_WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    # 中央配图
    img = IMG / "page-04-loom.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(2.29), Inches(1.20),
                             Inches(8.71), Inches(4.90))
    v1_bottom_bar(s,
                  killer="★ 杀手锏:MCP 工具集成  ·  Skill 技能库(改一个 .st 文件,AI 行为就变)",
                  note="对话持久化(H2 + Flyway)  ·  流式输出  ·  思维链可观测  ·  Spring Boot 自动配置")


# ====================================================================
# P5 SQL 工坊 + MCP
# ====================================================================
def slide_5_sql_forge(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    v1_title(s, accent=JAVA_BLUE_DARK,
             title="SQL 工坊 + MCP · 4 starter + 6 大功能",
             sub="4 个 starter 按需组合,AI 隔离在 5 个受限工具之后")
    # 右上角徽章
    add_rect(s, 11.40, 0.40, 1.70, 0.40,
             fill=JAVA_BLUE_DARK, line_color=None, line_width_pt=0)
    add_text(s, 11.40, 0.40, 1.70, 0.40, "🧩 4 starter",
             font=FONT_CN, size=11, color=COOL_WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    # 4 个 starter 色块
    starters = [
        "sql-forge 基础 CRUD",
        "+ calcite 跨库联邦",
        "+ web Amis 低代码",
        "+ mcp 本次演示",
    ]
    for i, name in enumerate(starters):
        x = 0.64 + i * 3.05
        add_rect(s, x, 1.20, 2.90, 0.50,
                 fill=DEEP_PANEL, line_color=None, line_width_pt=0)
        add_text(s, x, 1.20, 2.90, 0.50, name,
                 font=FONT_CN, size=12, color=COOL_WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    img = IMG / "page-05-sql-forge.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(2.83), Inches(1.80),
                             Inches(7.64), Inches(4.30))
    v1_bottom_bar(s,
                  killer="★ 杀手锏:JSON CRUD 协议  ·  Calcite 跨库联邦  ·  MCP 5 个受限工具",
                  note="AI 隔离后,私有化部署,数据不出企业")


# ====================================================================
# P6 演示 1(6 步流程色块,v1 风格但去时间戳)
# ====================================================================
def slide_6_demo1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    v1_title(s, accent=ORANGE,
             title="演示 1:一句话出分析报告",
             sub="演示基于 oms 订单管理系统(只用了 SQL 工坊的 基础 + MCP + Web 三个 starter)")
    # 演示语框
    add_rect(s, 0.60, 1.15, 12.10, 0.55,
             fill=DEEPER, line_color=None, line_width_pt=0)
    add_text(s, 0.75, 1.18, 12.00, 0.50,
             "💬 演示语:订单系统各分类商品数,画柱状图,保存 HTML 报告",
             font=FONT_CN, size=14, color=COOL_WHITE, bold=True,
             align=PP_ALIGN.LEFT)
    # 占位区
    add_rect(s, 0.60, 1.85, 12.10, 4.20,
             fill=DEEPER, line_color=None, line_width_pt=0)
    add_text(s, 0.60, 2.05, 12.10, 0.40,
             "📹 录屏占位(录完后自动替换为真截图)— 画面 6 步分解:",
             font=FONT_CN, size=12, color=COOL_DIM,
             align=PP_ALIGN.LEFT)
    steps = [
        "① 输入",
        "② AI 推理",
        "③ 表格+柱状图",
        "④ 保存 HTML",
        "⑤ 打开预览",
        "⑥ 展示链接",
    ]
    for i, step in enumerate(steps):
        x = 0.87 + i * 1.95
        add_rect(s, x, 3.05, 1.85, 0.95,
                 fill=DEEP_PANEL, line_color=None, line_width_pt=0)
        add_text(s, x, 3.05, 1.85, 0.95, step,
                 font=FONT_CN, size=12, color=COOL_WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text(s, 0.60, 5.40, 12.10, 0.40,
             "ffmpeg -i videos/demo1/demo1-final.mp4 -ss 30 -frames:v 1 images/page-06-demo1-thumb.png",
             font=FONT_MONO, size=10, color=COOL_DIM,
             align=PP_ALIGN.LEFT)
    v1_bottom_bar(s,
                  killer="90 秒出可用分析报告  —  录屏即将开始",
                  note="演示用的是基于 SQL 工坊搭建的 oms 订单管理系统,但 SQL 工坊不绑死订单,任何业务系统都能用")


# ====================================================================
# P7 演示 2
# ====================================================================
def slide_7_demo2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    v1_title(s, accent=ORANGE,
             title="演示 2:一句话生成 CRUD",
             sub="演示基于同一套 oms 订单管理系统(用 SQL 工坊的低代码 + CRUD 能力)")
    add_rect(s, 0.60, 1.15, 12.10, 0.55,
             fill=DEEPER, line_color=None, line_width_pt=0)
    add_text(s, 0.75, 1.18, 12.00, 0.50,
             "💬 演示语 2 轮:① 帮我做一个订单明细(ORDER_ITEMS)的维护页   ② 没问题,直接生成",
             font=FONT_CN, size=13, color=COOL_WHITE, bold=True,
             align=PP_ALIGN.LEFT)
    add_rect(s, 0.60, 1.85, 12.10, 4.20,
             fill=DEEPER, line_color=None, line_width_pt=0)
    add_text(s, 0.60, 2.05, 12.10, 0.40,
             "📹 录屏占位(录完后自动替换为真截图)— 画面 6 步分解:",
             font=FONT_CN, size=12, color=COOL_DIM,
             align=PP_ALIGN.LEFT)
    steps = [
        "① 输入+推断",
        "② 确认生成",
        "③ 3 行输出",
        "④ 切 oms",
        "⑤ 4 个功能",
        "⑥ 切回聊天",
    ]
    for i, step in enumerate(steps):
        x = 0.87 + i * 1.95
        add_rect(s, x, 3.05, 1.85, 0.95,
                 fill=DEEP_PANEL, line_color=None, line_width_pt=0)
        add_text(s, x, 3.05, 1.85, 0.95, step,
                 font=FONT_CN, size=12, color=COOL_WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text(s, 0.60, 5.40, 12.10, 0.40,
             "ffmpeg -i videos/demo2/demo2-final.mp4 -ss 50 -frames:v 1 images/page-07-demo2-thumb.png",
             font=FONT_MONO, size=10, color=COOL_DIM,
             align=PP_ALIGN.LEFT)
    v1_bottom_bar(s,
                  killer="100 秒出可用管理页  —  录屏即将开始",
                  note="演示用的 oms 订单系统只是 JavaBrain 的一种组合示例,任何业务库都能这么玩")


# ====================================================================
# P8 实战数据 + 对比表(v1 双栏布局)
# ====================================================================
def slide_8_comparison(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    v1_title(s, accent=ORANGE,
             title="7 轮实战验证 + 3 维度对比",
             sub="0 漂移 / 0 错误链接 / 0 乱码 / 5/5 自检 / >80% 节省")
    # 左:实战数据
    add_text(s, 0.60, 1.15, 5.50, 0.45,
             "📊 实战数据(技术派)",
             font=FONT_CN, size=18, color=JAVA_BLUE_DARK, bold=True,
             align=PP_ALIGN.LEFT)
    add_rect(s, 0.60, 1.65, 5.50, 4.45,
             fill=DEEP_PANEL, line_color=None, line_width_pt=0)
    add_text(s, 0.85, 1.85, 5.00, 0.60, "7 张表",
             font=FONT_CN, size=24, color=JAVA_BLUE_DARK, bold=True)
    add_text(s, 0.85, 2.45, 5.00, 0.35, "简单表 / 字典表 / 自关联 / 多外键",
             font=FONT_CN, size=10, color=COOL_GRAY)
    add_text(s, 0.85, 2.90, 5.00, 0.50,
             "0 漂移  ·  0 错误链接  ·  0 中文乱码",
             font=FONT_CN, size=13, color=SUCCESS, bold=True)
    add_text(s, 0.85, 3.50, 5.00, 0.45,
             "5 / 5  AI 字段推断自检通过",
             font=FONT_CN, size=14, color=COOL_WHITE, bold=True)
    add_text(s, 0.85, 4.05, 5.00, 0.65, "> 80%",
             font=FONT_CN, size=30, color=WARN, bold=True)
    add_text(s, 0.85, 4.70, 5.00, 0.35, "节省 CRUD 开发时间",
             font=FONT_CN, size=11, color=COOL_GRAY)
    # 配图(左下角小装饰)
    img = IMG / "page-08-comparison.png"
    if img.exists():
        add_picture_with_alpha(s, img, 4.10, 5.40, 1.90, 0.85, 60)
    # 右:对比表
    add_text(s, 6.50, 1.15, 6.30, 0.45,
             "📋 横向对比(商业派)",
             font=FONT_CN, size=18, color=PURPLE, bold=True,
             align=PP_ALIGN.LEFT)
    headers = ["维度", "低代码", "ChatGPT+插件", "JavaBrain"]
    cols_fill = [DEEP_PANEL, T.RGBColor(0x4B, 0x55, 0x63), DANGER, SUCCESS]
    cols_w = [1.40, 1.40, 1.70, 1.80]
    x0 = 6.50
    for c, (h, fc, w) in enumerate(zip(headers, cols_fill, cols_w)):
        x = x0 + sum(cols_w[:c])
        add_rect(s, x, 1.65, w, 0.50, fill=fc, line_color=None, line_width_pt=0)
        add_text(s, x, 1.65, w, 0.50, h,
                 font=FONT_CN, size=11, color=COOL_WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    rows = [
        ("业务取数",   "模板查询",       "✅ 但数据外发", "✅ 90秒+私有化"),
        ("CRUD 页面",  "不懂业务",       "需配合代码",    "✅ AI 语义推断"),
        ("跨库 JOIN",  "❌",              "❌",            "✅ Calcite 联邦"),
        ("私有化",     "✅",              "❌",            "✅"),
        ("列名识别",   "n/a",            "❌ 经常错",     "✅ 强制读 DDL"),
    ]
    for r, row in enumerate(rows):
        y = 2.15 + r * 0.50
        for c, cell in enumerate(row):
            x = x0 + sum(cols_w[:c])
            w = cols_w[c]
            fill_c = DARKER
            text_c = COOL_GRAY
            bold = False
            if c == 0:
                fill_c, text_c, bold = DEEP_PANEL, COOL_WHITE, True
            if c == 3:
                fill_c, text_c, bold = DEEPER, SUCCESS, True
            add_rect(s, x, y, w, 0.50, fill=fill_c, line_color=None, line_width_pt=0)
            add_text(s, x, y, w, 0.50, cell,
                     font=FONT_CN, size=11, color=text_c, bold=bold,
                     align=PP_ALIGN.CENTER)
    v1_bottom_bar(s,
                  killer="JavaBrain 在「安全 + 智能 + 私有化」三角全胜",
                  note="横向对比低代码 / ChatGPT+插件 — 数据不出企业,语义真推断,跨库真 JOIN")


# ====================================================================
# P9 开源 + 路线图
# ====================================================================
def slide_9_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    v1_title(s, accent=GREEN,
             title="已开源 + 未来规划 · 路线图",
             sub="三个独立仓库,各自独立维护,按需组合")
    # 3 个仓库卡片
    repos = [
        ("灵梭",         "spring-ai-loom-agent", "🧩 独立可用 · 7 模块 · AI 编排",   PURPLE),
        ("SQL 工坊",   "sql-forge",             "🧩 独立可用 · 12 模块 · 数据管理", JAVA_BLUE_DARK),
        ("SQL 工坊 MCP", "sql-forge-mcp",         "🧩 独立可用 · AI ↔ DB 工具",        GREEN),
    ]
    for i, (name, repo, sub, c) in enumerate(repos):
        x = 0.67 + i * 4.05
        add_rect(s, x, 1.20, 3.90, 1.00,
                 fill=DEEP_PANEL, line_color=None, line_width_pt=0)
        add_text(s, x + 0.20, 1.28, 3.50, 0.35, name,
                 font=FONT_CN, size=14, color=c, bold=True,
                 align=PP_ALIGN.LEFT)
        add_text(s, x + 0.20, 1.60, 3.50, 0.30, repo,
                 font=FONT_MONO, size=10, color=COOL_GRAY,
                 align=PP_ALIGN.LEFT)
        add_text(s, x + 0.20, 1.85, 3.50, 0.30, sub,
                 font=FONT_CN, size=9, color=COOL_DIM,
                 align=PP_ALIGN.LEFT)
    # 中央路线图配图
    img = IMG / "page-09-roadmap.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(3.36), Inches(2.40),
                             Inches(6.58), Inches(3.70))
    # V1.0 当前节点高亮(警示金)
    add_text(s, 4.5, 5.60, 4.3, 0.40,
             "V1.0 当前:基础 CRUD + NL2SQL + HTML 报告",
             font=FONT_CN, size=12, color=WARN, bold=True,
             align=PP_ALIGN.CENTER)
    v1_bottom_bar(s,
                  killer="★ 路线图 V1.0 当前 → V1.1 多租户 → V1.2 移动端 → V1.3 工作流引擎",
                  note="三个仓库独立维护,演示只是其中一种组合示例")


# ====================================================================
# P10 结尾(顶部 3 句深底 + 中央配图 + 底部黑条金句)
# ====================================================================
def slide_10_ending(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    v1_page_bg(s)
    # 顶部 3 句话深色块
    add_rect(s, 0.00, 0.30, 13.33, 1.95,
             fill=DEEPER, line_color=None, line_width_pt=0)
    add_text(s, 0.80, 0.45, 11.70, 0.50,
             "1.  JavaBrain = 开箱即用 + 私有化 + 业务语义",
             font=FONT_CN, size=18, color=SUCCESS, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, 0.80, 0.95, 11.70, 0.50,
             "2.  一行命令启动,一个文件改 AI 行为",
             font=FONT_CN, size=18, color=COOL_WHITE, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, 0.80, 1.45, 11.70, 0.50,
             "3.  让每一个 Spring Boot 项目都自带 AI 能力",
             font=FONT_CN, size=18, color=COOL_WHITE, bold=True,
             align=PP_ALIGN.LEFT)
    # 中央配图
    img = IMG / "page-10-ending.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(3.63), Inches(2.45),
                             Inches(6.04), Inches(3.40))
    # 底部黑条 + 橙金句
    add_rect(s, 0.00, 5.95, 13.33, 1.55,
             fill=DEEPER, line_color=None, line_width_pt=0)
    add_text(s, 0.80, 6.05, 11.70, 0.70,
             "让 AI 不再是 Demo,让企业系统真正智能。",
             font=FONT_CN, size=28, color=SUCCESS, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, 0.80, 6.75, 11.70, 0.40,
             "JavaBrain —— 一个文件改 AI 行为",
             font=FONT_CN, size=15, color=COOL_WHITE,
             align=PP_ALIGN.LEFT)
    add_text(s, 0.80, 7.10, 11.70, 0.35,
             "感谢聆听 · 欢迎交流",
             font=FONT_CN, size=12, color=COOL_DIM,
             align=PP_ALIGN.LEFT)


# ====================================================================
# 主入口
# ====================================================================
def main():
    prs = new_presentation()
    slide_1_cover(prs)
    slide_2_pain(prs)
    slide_3_position(prs)
    slide_4_loom(prs)
    slide_5_sql_forge(prs)
    slide_6_demo1(prs)
    slide_7_demo2(prs)
    slide_8_comparison(prs)
    slide_9_roadmap(prs)
    slide_10_ending(prs)
    prs.save(str(OUT))
    print(f"[OK] saved: {OUT}  ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
