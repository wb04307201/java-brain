"""T4 验证:P5-a/P5-b 应该有 7 slides,各种动画。"""
import sys
from pathlib import Path
from html import unescape

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

PPTX = Path(__file__).resolve().parent.parent / "javabrain.pptx"


def test_total_pages():
    prs = Presentation(str(PPTX))
    assert len(prs.slides) >= 7, f"期望 ≥7 页,实际 {len(prs.slides)}"


def test_p5a_animation_count():
    """动画已全部删除。"""
    prs = Presentation(str(PPTX))
    p5a = prs.slides[5]
    xml = etree.tostring(p5a._element).decode()
    timing_count = xml.count("<p:timing")
    assert timing_count == 0, f"P5-a 动画已删除,期望 0,实际 {timing_count}"


def test_p5a_pulse_indefinite():
    """动画已删除 — 无 pulse_loop。"""
    prs = Presentation(str(PPTX))
    p5a = prs.slides[5]
    xml = etree.tostring(p5a._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p5a_key_text():
    prs = Presentation(str(PPTX))
    p5a = prs.slides[5]
    xml = unescape(etree.tostring(p5a._element).decode())
    for text in ["SQL Forge", "从 JDBC 散弹到 Calcite 联邦", "Calcite",
                 "1 套协议代替 5 套 ORM", "数据不出企业"]:
        assert text in xml, f"P5-a 缺失: {text}"


def test_p5b_animation_count():
    """动画已全部删除。"""
    prs = Presentation(str(PPTX))
    p5b = prs.slides[6]
    xml = etree.tostring(p5b._element).decode()
    timing_count = xml.count("<p:timing")
    assert timing_count == 0, f"P5-b 动画已删除,期望 0,实际 {timing_count}"


def test_p5b_3_chase_indefinite():
    """动画已删除 — 无 pulse_loop。"""
    prs = Presentation(str(PPTX))
    p5b = prs.slides[6]
    xml = etree.tostring(p5b._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p5b_chase_delay_staggered():
    """动画已删除 — chase delay 不再需要。"""
    prs = Presentation(str(PPTX))
    p5b = prs.slides[6]
    xml = etree.tostring(p5b._element).decode()
    assert 'repeatCount' not in xml, "P5-b 应无残留 repeatCount"


def test_p5b_4_starters_present():
    prs = Presentation(str(PPTX))
    p5b = prs.slides[6]
    xml = unescape(etree.tostring(p5b._element).decode())
    for name in ["sql-forge-starter", "sql-forge-calcite",
                 "sql-forge-web", "sql-forge-mcp"]:
        assert name in xml, f"P5-b 缺失 starter: {name}"


def test_p5b_6_functions_present():
    prs = Presentation(str(PPTX))
    p5b = prs.slides[6]
    xml = unescape(etree.tostring(p5b._element).decode())
    for text in ["JSON CRUD", "Calcite", "MCP", "Amis", "实体", "零直连"]:
        assert text in xml, f"P5-b 缺失功能: {text}"


if __name__ == "__main__":
    test_total_pages()
    print("✓ 7 pages")
    test_p5a_animation_count()
    print("✓ P5-a animations")
    test_p5a_pulse_indefinite()
    print("✓ P5-a pulse_loop")
    test_p5a_key_text()
    print("✓ P5-a key text")
    test_p5b_animation_count()
    print("✓ P5-b animations")
    test_p5b_3_chase_indefinite()
    print("✓ P5-b 3 chase")
    test_p5b_chase_delay_staggered()
    print("✓ P5-b chase staggered")
    test_p5b_4_starters_present()
    print("✓ P5-b 4 starters")
    test_p5b_6_functions_present()
    print("✓ P5-b 6 functions")
    print("\nT4 验证全部通过 (动画已删除)")
