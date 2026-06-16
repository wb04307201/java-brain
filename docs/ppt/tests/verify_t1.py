"""T1 验证:P1 封面应该有 1 slide + 5 个 timing 节点 + 关键元素。"""
import sys
from pathlib import Path

# Force UTF-8 stdout on Windows to avoid GBK encode errors
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

PPTX = Path(__file__).resolve().parent.parent / "javabrain.pptx"


def test_p1_slide_count():
    prs = Presentation(str(PPTX))
    # T1 检查 P1 封面页存在(总页数 >= 1)
    assert len(prs.slides) >= 1, "expected at least 1 slide, got %d" % len(prs.slides)


def test_p1_animations_count():
    prs = Presentation(str(PPTX))
    slide = prs.slides[0]  # P1 是第 1 页(索引 0)
    xml = etree.tostring(slide._element).decode()
    # 6 个 animEffect 节点(大脑 zoom_in + 标题 + 副标 + 2 组件 + 金钩 zoom_in 模拟打字机)
    anim_count = xml.count("<p:animEffect")
    assert anim_count == 6, "expected 6 animEffects, got %d" % anim_count
    # 同时保证每页只有 1 个 <p:timing> 元素(累积式 add_anim 正确性)
    timing_count = xml.count("<p:timing")
    assert timing_count == 1, "expected 1 timing, got %d" % timing_count


def test_p1_key_text_present():
    prs = Presentation(str(PPTX))
    slide = prs.slides[0]
    xml = etree.tostring(slide._element).decode("utf-8")
    # python-pptx 把中文以 numeric XML entity 形式存储
    # 用 html.unescape 比对
    import html
    decoded = html.unescape(xml)
    for text in ["JavaBrain", "让您的系统瞬间拥有", "思考", "执行", "智能大脑", "3 分钟接入 AI"]:
        assert text in decoded, "missing text: " + text


def test_p1_colors_present():
    prs = Presentation(str(PPTX))
    slide = prs.slides[0]
    xml = etree.tostring(slide._element).decode()
    # Java 蓝 #007396 / AI 紫 #6B46C1 / 语义绿 #10B981 / 金 #F59E0B
    for hex_color in ["007396", "6B46C1", "10B981", "F59E0B"]:
        assert hex_color in xml, "missing color: #" + hex_color


if __name__ == "__main__":
    test_p1_slide_count()
    print("[OK] 1 page")
    test_p1_animations_count()
    print("[OK] 6 animEffects + 1 timing")
    test_p1_key_text_present()
    print("[OK] key text")
    test_p1_colors_present()
    print("[OK] 4 colors")
    print("\nT1 all checks passed")
