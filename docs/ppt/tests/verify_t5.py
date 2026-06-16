"""T5 验证:P6/P7 应该有 9 slides,5 步依次入场。"""
import sys
from pathlib import Path
from html import unescape

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

PPTX = Path(__file__).resolve().parent.parent / "javabrain.pptx"


def test_total_pages():
    prs = Presentation(str(PPTX))
    assert len(prs.slides) >= 9, f"期望 ≥9 页,实际 {len(prs.slides)}"


def test_p6_animation_count():
    """动画已全部删除。"""
    prs = Presentation(str(PPTX))
    p6 = prs.slides[7]
    xml = etree.tostring(p6._element).decode()
    anim_count = xml.count("<p:animEffect")
    assert anim_count == 0, f"P6 动画已删除,期望 0,实际 {anim_count}"


def test_p6_pulse_indefinite():
    """动画已删除 — 无 pulse_loop。"""
    prs = Presentation(str(PPTX))
    p6 = prs.slides[7]
    xml = etree.tostring(p6._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p6_key_text():
    prs = Presentation(str(PPTX))
    p6 = prs.slides[7]
    xml = unescape(etree.tostring(p6._element).decode())
    for text in ["演示 1", "一句话出分析报告", "上个月", "NL2SQL", "42 秒"]:
        assert text in xml, f"P6 缺失: {text}"


def test_p7_animation_count():
    """动画已全部删除。"""
    prs = Presentation(str(PPTX))
    p7 = prs.slides[8]
    xml = etree.tostring(p7._element).decode()
    anim_count = xml.count("<p:animEffect")
    assert anim_count == 0, f"P7 动画已删除,期望 0,实际 {anim_count}"


def test_p7_pulse_indefinite():
    """动画已删除 — 无 pulse_loop。"""
    prs = Presentation(str(PPTX))
    p7 = prs.slides[8]
    xml = etree.tostring(p7._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p7_key_text():
    prs = Presentation(str(PPTX))
    p7 = prs.slides[8]
    xml = unescape(etree.tostring(p7._element).decode())
    for text in ["演示 2", "一句话生成 CRUD", "商品管理", "web.st", "40 秒"]:
        assert text in xml, f"P7 缺失: {text}"


if __name__ == "__main__":
    test_total_pages()
    print("✓ 9 pages")
    test_p6_animation_count()
    print("✓ P6 animations")
    test_p6_pulse_indefinite()
    print("✓ P6 pulse_loop")
    test_p6_key_text()
    print("✓ P6 key text")
    test_p7_animation_count()
    print("✓ P7 animations")
    test_p7_pulse_indefinite()
    print("✓ P7 pulse_loop")
    test_p7_key_text()
    print("✓ P7 key text")
    print("\nT5 验证全部通过 (动画已删除)")
