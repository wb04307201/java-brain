"""T7 验证:12 页完成,P9+P10 关键内容。

累计动画统计(P9 重做为'未来目标'页后):
- P1: 5 fade_in (no pulse)
- P2: 7 fade_in + 1 pulse_loop = 8
- P3: 3 fade_in + 1 pulse_loop = 4
- P4-a: 3 fade_in + 1 pulse_loop = 4
- P4-b: 6 fade_in + 2 pulse_loop (chase) = 8
- P5-a: 3 fade_in + 1 pulse_loop = 4
- P5-b: 10 fade_in + 3 pulse_loop (chase) = 13
- P6: 5 fade_in + 1 pulse_loop = 6
- P7: 5 fade_in + 1 pulse_loop = 6
- P8: 9 fade_in + 1 pulse_loop = 10
- P9: 12 fade_in + 1 pulse_loop (CDA 金钩) = 13  ← 新增未来目标,内容更密
- P10: 5 fade_in + 1 pulse_loop (killer) = 6

合计:73 fade_in + 14 pulse_loop = 87 anim,14 indefinite。
(P9 重做前是 79,新增 8 个 fade_in 用于灵梭/SQL工坊各 4 张卡入场)
"""
import sys
from pathlib import Path
from html import unescape

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

PPTX = Path(__file__).resolve().parent.parent / "javabrain.pptx"


def test_total_pages():
    prs = Presentation(str(PPTX))
    assert len(prs.slides) == 12, f"期望 12 页,实际 {len(prs.slides)}"


def test_total_animations():
    """12 页累计动画总数(P1 加大脑 zoom_in + 金钩 zoom_in 模拟打字机 + P10 logo spin)。

    按实际 add_anim 调用累计:
    - 85 fade_in/zoom_in/spin + 15 pulse_loop(其中 5 个 chase)= 100 anim,15 indefinite
    """
    prs = Presentation(str(PPTX))
    total = 0
    total_indefinite = 0
    for slide in prs.slides:
        xml = etree.tostring(slide._element).decode()
        # 改用 animEffect 计数(累积式 add_anim 修复后每页只有 1 个 <p:timing>)
        total += xml.count("<p:animEffect")
        total_indefinite += xml.count('repeatCount="indefinite"')
    print(f"  total animEffects: {total}, indefinites: {total_indefinite}")
    assert total == 100, f"期望 100 animEffect,实际 {total}"
    assert total_indefinite == 15, f"期望 15 pulse_loop,实际 {total_indefinite}"


def test_p9_loom_future():
    """P9 灵梭 4 未来目标(Agentic Workflow / Graph-RAG / 多模态 / Skill 市场)。"""
    prs = Presentation(str(PPTX))
    p9 = prs.slides[10]
    xml = unescape(etree.tostring(p9._element).decode())
    for text in ["灵梭", "Agentic Workflow", "Graph-RAG",
                 "多模态", "Skill"]:
        assert text in xml, f"P9 灵梭缺失: {text}"


def test_p9_forge_future():
    """P9 SQL工坊 4 未来目标(Text-to-SQL / HTAP / 数据治理 / Serverless)。"""
    prs = Presentation(str(PPTX))
    p9 = prs.slides[10]
    xml = unescape(etree.tostring(p9._element).decode())
    for text in ["SQL工坊", "Text-to-SQL", "HTAP",
                 "Serverless", "数据治理"]:
        assert text in xml, f"P9 SQL工坊缺失: {text}"


def test_p9_dual_engine_vision():
    """P9 底部双引擎协同 3 金钩。"""
    prs = Presentation(str(PPTX))
    p9 = prs.slides[10]
    xml = unescape(etree.tostring(p9._element).decode())
    for text in ["AI 懂数据", "数据懂 AI", "CDA",
                 "零代码 AI 应用工厂"]:
        assert text in xml, f"P9 双引擎金钩缺失: {text}"


def test_p9_pulse_indefinite():
    """CDA 金钩脉冲。"""
    prs = Presentation(str(PPTX))
    p9 = prs.slides[10]
    xml = etree.tostring(p9._element).decode()
    assert 'repeatCount="indefinite"' in xml


def test_p9_animation_count():
    """P9:13 动画(1 标题 + 4 灵梭 + 4 SQL工坊 + 3 金钩 fade + 1 金钩 pulse)。"""
    prs = Presentation(str(PPTX))
    p9 = prs.slides[10]
    xml = etree.tostring(p9._element).decode()
    tc = xml.count("<p:animEffect")
    assert tc == 13, f"P9 期望 13 animEffect,实际 {tc}"


def test_p10_key_text():
    prs = Presentation(str(PPTX))
    p10 = prs.slides[11]
    xml = unescape(etree.tostring(p10._element).decode())
    for text in ["JavaBrain", "开箱即用", "AI 懂数据",
                 "LoomAgentApplication", "25.394",
                 "github.com/wb04307201"]:
        assert text in xml, f"P10 缺失: {text}"


def test_p10_three_sentences():
    """3 句话:开箱即用 / 一行命令 / 每个 Spring Boot。"""
    prs = Presentation(str(PPTX))
    p10 = prs.slides[11]
    xml = unescape(etree.tostring(p10._element).decode())
    for txt in ["开箱即用", "一行命令启动", "每个 Spring Boot"]:
        assert txt in xml, f"P10 缺失: {txt}"


def test_p10_pulse_indefinite():
    """金句持续脉冲。"""
    prs = Presentation(str(PPTX))
    p10 = prs.slides[11]
    xml = etree.tostring(p10._element).decode()
    indefinite_count = xml.count('repeatCount="indefinite"')
    assert indefinite_count >= 1, f"P10 期望 ≥1 pulse_loop(金句),实际 {indefinite_count}"


def test_p10_animation_count():
    """P10:7 动画(logo spin + title fade + 3 sent fade + killer pulse + term fade)。"""
    prs = Presentation(str(PPTX))
    p10 = prs.slides[11]
    xml = etree.tostring(p10._element).decode()
    tc = xml.count("<p:animEffect")
    assert tc == 7, f"P10 期望 7 animEffect,实际 {tc}"


if __name__ == "__main__":
    test_total_pages()
    print("OK 12 pages")
    test_total_animations()
    print("OK 100 animEffects + 15 indefinite")
    test_p9_loom_future()
    print("OK P9 灵梭 4 future")
    test_p9_forge_future()
    print("OK P9 SQL工坊 4 future")
    test_p9_dual_engine_vision()
    print("OK P9 双引擎金钩")
    test_p9_pulse_indefinite()
    print("OK P9 CDA pulse")
    test_p9_animation_count()
    print("OK P9 13 anim")
    test_p10_key_text()
    print("OK P10 key text")
    test_p10_three_sentences()
    print("OK P10 3 sentences")
    test_p10_pulse_indefinite()
    print("OK P10 pulse_loop")
    test_p10_animation_count()
    print("OK P10 7 anim")
    print("\nT7 验证全部通过 + 12 页 PPT 完成")
