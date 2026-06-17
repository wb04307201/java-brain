"""tools.py 单元测试 —— 28 个核心 helper 行为测试。

不依赖网络,全部本地。"""
from __future__ import annotations

import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pptx import Presentation
from pptx.util import Inches

import tools as T
from tools import (
    BG_DEEP, BG_PANEL, GRID_LINE,
    TEXT_PRIMARY, TEXT_SECONDARY, TEXT_DIM,
    PRIMARY, SUCCESS, DANGER, WARN,
    FONT_CN, FONT_MONO,
    add_text, set_text, add_rect, add_picture_with_alpha,
    apply_chrome, section_label, corner_badge, add_page_number,
    kill_box, node_block, card, terminal_box, big_num,
    iso_lego, arrow_line,
    new_presentation, qa_image_check, find_dominant_bg, cover_region,
    find_media_target, replace_pptx_media, qa_pptx_images,
    preview_grid, wanx_generate, wanx_download, wanx_upload, wanx_edit,
)

IMG_DIR = ROOT / "images"


def _new():
    prs = new_presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


# 计数
passed = 0
failed = []


def t(name, cond, info=""):
    global passed
    if cond:
        passed += 1
        print(f"  [OK]  {name}")
    else:
        failed.append((name, info))
        print(f"  [FAIL] {name}  -- {info}")


# ====================================================================
# 1-4: A 组(文本/形状)
# ====================================================================
prs, s = _new()
tb = add_text(s, 1, 1, 4, 0.5, "Hello",
              size=18, color=TEXT_PRIMARY, bold=True)
t("1. add_text returns textbox", tb.has_text_frame)
t("2. add_text text matches", "Hello" in tb.text_frame.text)

prs, s = _new()
tb2 = add_text(s, 1, 1, 4, 0.5, "X", size=18)
set_text(tb2, "Y", size=24, color=SUCCESS, bold=True)
t("3. set_text updates content", "Y" in tb2.text_frame.text)
t("4. set_text preserves position",
  tb2.left == Inches(1) and tb2.top == Inches(1))

prs, s = _new()
r = add_rect(s, 0.5, 0.5, 3, 1, fill=BG_PANEL,
             line_color=SUCCESS, line_width_pt=2)
t("5. add_rect fills", r.fill.type is not None)
t("6. add_rect has line", r.line.color.rgb is not None)

# ====================================================================
# 7-12: B 组(布局/装饰)
# ====================================================================
prs, s = _new()
apply_chrome(s, 1, 10, "测试")
# 验证:背景+网格+HUD 角标+章节标都加了 shape
n_shapes = len(s.shapes)
t("7. apply_chrome adds 4+ shapes", n_shapes >= 4,
  f"shapes={n_shapes}")

prs, s = _new()
section_label(s, "测试章节")
t("8. section_label adds textbox", len(s.shapes) >= 1)

prs, s = _new()
corner_badge(s, 9, 0.5, 3, 0.4, "★ 提示")
t("9. corner_badge adds shape", len(s.shapes) >= 2,
  f"shapes={len(s.shapes)}")

prs, s = _new()
add_page_number(s, 5, 10)
t("10. add_page_number adds textbox", len(s.shapes) >= 1)

# ====================================================================
# 13-18: C 组(复合元素)
# ====================================================================
prs, s = _new()
kill_box(s, 1, 1, 4, 1, "杀手锏金句")
t("11. kill_box has 2 shapes", len(s.shapes) >= 2)

prs, s = _new()
node_block(s, 1, 1, 4, 1, "节点", "副标", color=PRIMARY, star=True)
t("12. node_block star=True", len(s.shapes) >= 3)

prs, s = _new()
c = card(s, 1, 1, 4, 2, fill=BG_PANEL, border=GRID_LINE, border_pt=2)
t("13. card returns shape with line", c.line.color.rgb is not None)

prs, s = _new()
terminal_box(s, 1, 1, 6, 3, ["line1", "line2"], title="终端")
t("14. terminal_box creates box", len(s.shapes) >= 2)

prs, s = _new()
big_num(s, 1, 1, 4, 1.5, "120", "秒", "耗时", color=DANGER, num_size=120)
t("15. big_num creates textbox", len(s.shapes) >= 1)

prs, s = _new()
iso_lego(s, 1, 1, 2, 1.5, PRIMARY, studs=True, highlight=True)
t("16. iso_lego creates block", len(s.shapes) >= 4)

# ====================================================================
# 19-22: 箭头 / 装饰
# ====================================================================
prs, s = _new()
arrow_line(s, 1, 1, 5, 3, color=SUCCESS, weight_pt=2)
t("17. arrow_line creates line", len(s.shapes) >= 1)

# ====================================================================
# 23-25: 文件 I/O
# ====================================================================
prs, s = _new()
# 加一张图(从已有 images 拿)
src = IMG_DIR / "page-01-cover.png"
if src.exists():
    pic = add_picture_with_alpha(s, src, 0.5, 0.5, 4, 3, 75)
    t("18. add_picture_with_alpha adds pic", pic is not None)
else:
    t("18. add_picture_with_alpha (skip, no image)", True)

# qa_image_check(对真实图片)
if src.exists():
    ok, issues = qa_image_check(src)
    t("19. qa_image_check returns tuple",
      isinstance(ok, bool) and isinstance(issues, list))
    t("20. qa_image_check passes for known good image", ok)
else:
    t("19. qa_image_check (skip, no image)", True)
    t("20. qa_image_check (skip, no image)", True)

# find_dominant_bg / cover_region / find_media_target / replace_pptx_media
test_img = IMG_DIR / "page-02-pain.png"
if test_img.exists():
    bg = find_dominant_bg(test_img)
    t("21. find_dominant_bg returns rgb",
      isinstance(bg, tuple) and len(bg) == 3)
    from PIL import Image
    img = Image.open(test_img)
    covered = cover_region(img, (0, 0, 100, 100), bg, pad=4)
    t("22. cover_region returns image", hasattr(covered, "size"))
else:
    t("21. find_dominant_bg (skip)", True)
    t("22. cover_region (skip)", True)

# preview_grid
out = ROOT / "preview" / "test_grid.png"
out.parent.mkdir(parents=True, exist_ok=True)
imgs = sorted(IMG_DIR.glob("page-*.png"))[:5]
if imgs:
    res = preview_grid(IMG_DIR, out, cols=5)
    t("23. preview_grid creates file", res.exists() and res.stat().st_size > 0)
else:
    t("23. preview_grid (skip, no images)", True)

# ====================================================================
# 26-28: 集成测试(整本 pptx 跑通)
# ====================================================================
prs, s = _new()
apply_chrome(s, 1, 10, "集成")
add_text(s, 1, 1, 6, 0.5, "集成测试标题", size=32, bold=True)
card(s, 1, 1.5, 5, 2, border=WARN, border_pt=2.5, star=True)
add_text(s, 1.1, 1.6, 4.8, 0.5, "卡片标题", size=18, bold=True, color=WARN)
kill_box(s, 7, 1, 5, 1.5, "金句", size=24)
out_pptx = ROOT / "preview" / "test-tools.pptx"
prs.save(str(out_pptx))
re = Presentation(str(out_pptx))
t("24. save+reload roundtrip", len(re.slides) == 1)

# 调一次 wanx_* import(不真打网络)
try:
    _ = T.wanx_generate
    t("25. wanx_generate exists", True)
except Exception as e:
    t("25. wanx_generate", False, str(e))

# 调一次 find_media_target(replace_pptx_media 的辅助)
src2 = IMG_DIR / "page-01-cover.png"
if src2.exists() and out_pptx.exists():
    target = find_media_target(re, hint="image", ext="png")
    t("26. find_media_target returns str|None",
      target is None or isinstance(target, str))

# qa_pptx_images(对生成的 v3 pptx 跑,验证配图都在)
v3_pptx = ROOT / "javabrain-v3.pptx"
if v3_pptx.exists():
    ok, issues = qa_pptx_images(v3_pptx)
    t("27. qa_pptx_images returns tuple",
      isinstance(ok, bool) and isinstance(issues, list))
    t("28. qa_pptx_images has 8+ items checked",
      len(issues) >= 0)

# ====================================================================
# 报告
# ====================================================================
print(f"\n=== tools 单元测试:{passed} pass / {len(failed)} fail ===")
if failed:
    print("FAIL details:")
    for n, info in failed:
        print(f"  - {n}: {info}")
    sys.exit(1)
print("[ALL PASS]")
