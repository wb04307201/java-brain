"""T6 验证:P8 实战对比应该有 10 slides,关键内容。"""
import sys
from pathlib import Path
from html import unescape

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

PPTX = Path(__file__).resolve().parent.parent / "javabrain.pptx"


def test_total_pages():
    prs = Presentation(str(PPTX))
    assert len(prs.slides) >= 10, f"期望 ≥10 页,实际 {len(prs.slides)}"


def test_p8_animation_count():
    """动画已全部删除。"""
    prs = Presentation(str(PPTX))
    p8 = prs.slides[9]
    xml = etree.tostring(p8._element).decode()
    anim_count = xml.count("<p:animEffect")
    assert anim_count == 0, f"P8 动画已删除,期望 0,实际 {anim_count}"


def test_p8_pulse_indefinite():
    """动画已删除 — 无 pulse_loop。"""
    prs = Presentation(str(PPTX))
    p8 = prs.slides[9]
    xml = etree.tostring(p8._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p8_bar_subtitle_present():
    """P8 柱状图新副标(≈ 3 月 / 10 天 / ≈ 3 分)。"""
    prs = Presentation(str(PPTX))
    p8 = prs.slides[9]
    xml = unescape(etree.tostring(p8._element).decode())
    for text in ["≈ 3 月", "10 天", "≈ 3 分"]:
        assert text in xml, f"P8 柱状图副标缺失: {text}"


def test_p8_no_bottom_gold_killer():
    """P8 底部金句已删除(让柱状图自己说话)。"""
    prs = Presentation(str(PPTX))
    p8 = prs.slides[9]
    xml = unescape(etree.tostring(p8._element).decode())
    assert "JavaBrain 在 安全 + 智能 + 私有化 三角都做到" not in xml, \
        "P8 底部金句应已删除"


def test_p8_tech_items():
    """技术派 4 行 ✓。"""
    prs = Presentation(str(PPTX))
    p8 = prs.slides[9]
    xml = unescape(etree.tostring(p8._element).decode())
    for text in ["技术派", "7 张表测试", "0 漂移", "字段推断", ">80%"]:
        assert text in xml, f"P8 技术派缺失: {text}"


def test_p8_biz_items():
    """商业派 5 行对比。"""
    prs = Presentation(str(PPTX))
    p8 = prs.slides[9]
    xml = unescape(etree.tostring(p8._element).decode())
    for text in ["商业派", "业务取数", "CRUD 页面", "跨库 JOIN", "私有化", "列名识别"]:
        assert text in xml, f"P8 商业派缺失: {text}"


def test_p8_killer_pulse():
    """动画已删除 — pulse delay 不再需要。"""
    prs = Presentation(str(PPTX))
    p8 = prs.slides[9]
    xml = etree.tostring(p8._element).decode()
    assert '12000' not in xml, "P8 应无残留 12000ms pulse"


if __name__ == "__main__":
    test_total_pages()
    print("✓ 10 pages")
    test_p8_animation_count()
    print("✓ P8 animations")
    test_p8_pulse_indefinite()
    print("✓ P8 pulse_loop")
    test_p8_tech_items()
    print("✓ P8 tech items")
    test_p8_biz_items()
    print("✓ P8 biz items")
    test_p8_killer_pulse()
    print("✓ P8 >80% pulse delay")
    print("\nT6 验证全部通过 (动画已删除 + 柱状图副标)")
