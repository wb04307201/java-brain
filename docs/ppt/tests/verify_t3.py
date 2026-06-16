"""T3 验证:P3/P4-a/P4-b 应该有 5 slides + 各种动画数。"""
import sys
from pathlib import Path
from html import unescape

# Force UTF-8 stdout on Windows to avoid GBK encode errors
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

PPTX = Path(__file__).resolve().parent.parent / "javabrain.pptx"


def test_total_pages():
    prs = Presentation(str(PPTX))
    assert len(prs.slides) >= 5, f"期望 ≥5 页,实际 {len(prs.slides)}"


def test_p3_animation_count():
    """动画已全部删除 — 用户手动添加。"""
    prs = Presentation(str(PPTX))
    p3 = prs.slides[2]
    xml = etree.tostring(p3._element).decode()
    anim_count = xml.count("<p:animEffect")
    assert anim_count == 0, f"P3 动画已删除,期望 0,实际 {anim_count}"


def test_p3_pulse_indefinite():
    """动画已删除 — 无 pulse_loop。"""
    prs = Presentation(str(PPTX))
    p3 = prs.slides[2]
    xml = etree.tostring(p3._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p3_key_text():
    """P3 4 行杀手锏:灵梭 / SQL工坊 / SQL工坊 MCP / 全部杀手锏。"""
    prs = Presentation(str(PPTX))
    p3 = prs.slides[2]
    xml = unescape(etree.tostring(p3._element).decode())
    for text in ["JavaBrain", "灵梭杀手锏", "SQL工坊杀手锏", "SQL工坊 MCP 杀手锏",
                 "全部杀手锏", "数据不出企业", "依赖库"]:
        assert text in xml, f"P3 缺失: {text}"


def test_p4a_animation_count():
    """动画已全部删除。"""
    prs = Presentation(str(PPTX))
    p4a = prs.slides[3]
    xml = etree.tostring(p4a._element).decode()
    anim_count = xml.count("<p:animEffect")
    assert anim_count == 0, f"P4-a 动画已删除,期望 0,实际 {anim_count}"


def test_p4a_pulse_indefinite():
    """动画已删除 — 无 pulse_loop。"""
    prs = Presentation(str(PPTX))
    p4a = prs.slides[3]
    xml = etree.tostring(p4a._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p4a_key_text():
    prs = Presentation(str(PPTX))
    p4a = prs.slides[3]
    xml = unescape(etree.tostring(p4a._element).decode())
    # 新金句(P4-a 方案 A:数字型)
    for text in ["LoomAgent", "从 Spring AI 裸用到 JavaBrain 一体化",
                 "灵梭 starter", ".st 模板", "50 行配置",
                 "1 行模板", "49 行省下"]:
        assert text in xml, f"P4-a 缺失: {text}"


def test_p4b_animation_count():
    """动画已全部删除。"""
    prs = Presentation(str(PPTX))
    p4b = prs.slides[4]
    xml = etree.tostring(p4b._element).decode()
    anim_count = xml.count("<p:animEffect")
    assert anim_count == 0, f"P4-b 动画已删除,期望 0,实际 {anim_count}"


def test_p4b_chase_indefinite():
    """动画已删除 — 无 pulse_loop。"""
    prs = Presentation(str(PPTX))
    p4b = prs.slides[4]
    xml = etree.tostring(p4b._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p4b_chase_delay_staggered():
    """动画已删除 — chase delay 不再需要,改为检查无 delay 残留。"""
    prs = Presentation(str(PPTX))
    p4b = prs.slides[4]
    xml = etree.tostring(p4b._element).decode()
    assert 'repeatCount' not in xml, "P4-b 应无残留 repeatCount"


def test_p4b_modules_present():
    prs = Presentation(str(PPTX))
    p4b = prs.slides[4]
    xml = unescape(etree.tostring(p4b._element).decode())
    # Row1 杀手锏(RAG★/MCP★/Skill★) + Row2 基础(对话/文件/工具)
    for text in ["RAG 知识库", "MCP", "Skill", "文件管理", "对话 UI", "内置工具"]:
        assert text in xml, f"P4-b 缺失模块: {text}"


if __name__ == "__main__":
    test_total_pages()
    print("[OK] 5 pages")
    test_p3_animation_count()
    print("[OK] P3 animations")
    test_p3_pulse_indefinite()
    print("[OK] P3 pulse_loop")
    test_p3_key_text()
    print("[OK] P3 key text")
    test_p4a_animation_count()
    print("[OK] P4-a animations")
    test_p4a_pulse_indefinite()
    print("[OK] P4-a pulse_loop")
    test_p4a_key_text()
    print("[OK] P4-a key text")
    test_p4b_animation_count()
    print("[OK] P4-b animations")
    test_p4b_chase_indefinite()
    print("[OK] P4-b 2 chase")
    test_p4b_chase_delay_staggered()
    print("[OK] P4-b chase staggered 2s")
    test_p4b_modules_present()
    print("[OK] P4-b 6 modules")
    print("\nT3 验证全部通过 (动画已删除)")
