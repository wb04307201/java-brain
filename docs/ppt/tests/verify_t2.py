"""T2 验证:P2 痛点应该有 2 slides + 第 2 页有 8 个动画(7 fade_in + 1 pulse_loop indefinite)+ 关键元素。"""
import sys
from pathlib import Path
from html import unescape

# Force UTF-8 stdout on Windows to avoid GBK encode errors
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

PPTX = Path(__file__).resolve().parent.parent / "javabrain.pptx"


def test_total_pages():
    prs = Presentation(str(PPTX))
    assert len(prs.slides) >= 2, f"期望 ≥2 页,实际 {len(prs.slides)}"


def test_p2_animation_count():
    """动画已全部删除 — 用户手动添加。"""
    prs = Presentation(str(PPTX))
    p2 = prs.slides[1]
    xml = etree.tostring(p2._element).decode()
    timing_count = xml.count("<p:timing")
    assert timing_count == 0, f"无 timing,期望 0,实际 {timing_count}"


def test_p2_pulse_loop_indefinite():
    """动画已删除 — pulse_loop 也归零。"""
    prs = Presentation(str(PPTX))
    p2 = prs.slides[1]
    xml = etree.tostring(p2._element).decode()
    assert 'repeatCount="indefinite"' not in xml


def test_p2_key_text():
    prs = Presentation(str(PPTX))
    p2 = prs.slides[1]
    xml = unescape(etree.tostring(p2._element).decode())
    for text in ["3 月", "5 天", "3 天", "3 分", "90 秒", "10 分",
                 "1440×", "接入 AI"]:
        assert text in xml, f"缺失文字: {text}"


def test_p2_red_green_gold_colors():
    prs = Presentation(str(PPTX))
    p2 = prs.slides[1]
    xml = etree.tostring(p2._element).decode()
    for hex_color in ["DC2626", "10B981", "F59E0B"]:
        assert hex_color in xml, f"缺失颜色: #{hex_color}"


if __name__ == "__main__":
    test_total_pages()
    print("[OK] 2 pages")
    test_p2_animation_count()
    print("[OK] 0 animations (动画已删除)")
    test_p2_pulse_loop_indefinite()
    print("[OK] 0 pulse_loop")
    test_p2_key_text()
    print("[OK] key text")
    test_p2_red_green_gold_colors()
    print("[OK] red/green/gold colors")
    print("\nT2 all checks passed")
