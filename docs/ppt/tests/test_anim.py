"""add_anim() 工具函数单元测试。

通过把 gen-ppt.py 当脚本运行(sys.path insert)导入 gen_ppt 模块,
验证 add_anim() 在 PowerPoint XML 中正确生成:
1. fade_in 生成 <p:timing> 节点 + presetClass="entr" presetId="10"
2. pulse_loop 生成 repeatCount="indefinite"
3. delay_ms 反映到 begin 属性(2000ms → XML 含 '2000')
"""
import importlib.util
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))
# 文件名带连字符(gen-ppt.py)不能直接 import,用 importlib 加载
_spec = importlib.util.spec_from_file_location(
    "gen_ppt",
    Path(__file__).resolve().parent.parent / "scripts" / "gen-ppt.py",
)
gen_ppt = importlib.util.module_from_spec(_spec)  # type: ignore
_spec.loader.exec_module(gen_ppt)  # type: ignore

from lxml import etree  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402


def test_fade_in_creates_timing_node():
    """fade_in 应在 slide 上生成 <p:timing> 节点 + entr/10 preset。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
    gen_ppt.add_anim(slide, rect, "fade_in", delay_ms=0, dur_ms=500, loop=False)
    xml = etree.tostring(slide._element).decode()
    assert "<p:timing" in xml
    assert 'presetClass="entr"' in xml
    assert 'presetId="10"' in xml


def test_pulse_loop_has_indefinite():
    """pulse_loop 应生成 repeatCount='indefinite'(持续到翻页,转视频核心)。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
    gen_ppt.add_anim(slide, rect, "pulse", delay_ms=0, dur_ms=1500, loop=True)
    xml = etree.tostring(slide._element).decode()
    assert 'repeatCount="indefinite"' in xml
    # pulse 是 emph 类的 26 号
    assert 'presetClass="emph"' in xml
    assert 'presetId="26"' in xml


def test_chase_pulse_with_delay():
    """delay_ms 错开应生成对应的 begin 属性(chase_pulse 基础)。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
    gen_ppt.add_anim(slide, rect, "pulse", delay_ms=2000, dur_ms=1500, loop=True)
    xml = etree.tostring(slide._element).decode()
    assert '2000' in xml  # begin="2000ms"
